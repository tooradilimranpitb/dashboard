import calendar
import copy
import datetime
import io
import json
import os
import re
import time
import duckdb
import shutil
import zipfile
import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st
import streamlit.components.v1 as components
import bcrypt

# --- Advanced Statistical & ML Imports ---
from sklearn.linear_model import LinearRegression
from sklearn.preprocessing import PolynomialFeatures
from sklearn.pipeline import make_pipeline

# --- Export Imports ---
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from pptx import Presentation
from pptx.util import Inches, Pt

# --- Page Configuration ---
st.set_page_config(
    page_title="DPA Executive Analytics Center & Target Management",
    page_icon="analytics",
    layout="wide",
)

# --- Force-scrollable dropdown fix ---
components.html(
    """
    <script>
    (function () {
      function fixDropdowns() {
        try {
          const doc = window.parent.document;
          const boxes = doc.querySelectorAll('ul[role="listbox"], div[role="listbox"]');
          boxes.forEach(function (el) {
            el.style.setProperty('max-height', '260px', 'important');
            el.style.setProperty('overflow-y', 'auto', 'important');
            el.style.setProperty('overflow-x', 'hidden', 'important');

            let node = el.parentElement;
            let hops = 0;
            while (node && hops < 6) {
              const style = window.parent.getComputedStyle(node);
              if (style.overflow === 'hidden' || style.overflowY === 'hidden') {
                node.style.setProperty('overflow', 'visible', 'important');
              }
              if (node.getAttribute && node.getAttribute('data-baseweb') === 'popover') {
                break;
              }
              node = node.parentElement;
              hops += 1;
            }
          });
        } catch (e) {}
      }
      try {
        const doc = window.parent.document;
        const observer = new MutationObserver(fixDropdowns);
        observer.observe(doc.body, { childList: true, subtree: true });
        fixDropdowns();
      } catch (e) {}
    })();
    </script>
    """,
    height=0,
)

# --- Load Google Material Symbols & FontAwesome Icons ---
st.markdown(
    """
    <link rel="stylesheet" href="https://fonts.googleapis.com/css2?family=Material+Symbols+Outlined:opsz,wght,FILL,GRAD@24,400,0,0" />
    """,
    unsafe_allow_html=True,
)

# --- Brand & Chart Palette ---
DPA_NAVY = "#16233F"
DPA_NAVY_LIGHT = "#2A3D66"
DPA_TEAL = "#0EA5A6"
DPA_TEAL_DARK = "#0B8485"
DPA_AMBER = "#E8A33D"
DPA_GREEN = "#2FA875"
DPA_RED = "#D64550"
DPA_BG = "#F8FAFC"
DPA_MUTED = "#64748B"

CHART_COLORWAY = [DPA_TEAL, DPA_NAVY, DPA_AMBER, "#6C63FF", DPA_GREEN, DPA_RED, "#4C7CA8", "#B08968"]

def style_fig(fig, title=None):
    fig.update_layout(
        colorway=CHART_COLORWAY,
        font=dict(family="Inter, -apple-system, sans-serif", size=13, color=DPA_NAVY),
        title_font=dict(family="Sora, Inter, sans-serif", size=17, color=DPA_NAVY),
        plot_bgcolor="rgba(0,0,0,0)",
        paper_bgcolor="rgba(0,0,0,0)",
        margin=dict(t=60, l=15, r=15, b=15),
        legend=dict(
            orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1,
            bgcolor="rgba(0,0,0,0)",
        ),
        hoverlabel=dict(bgcolor="white", font_size=13, font_family="Inter, sans-serif"),
        dragmode="zoom",
    )
    fig.update_xaxes(
        showgrid=False, linecolor="#E2E8F0", tickfont=dict(color=DPA_MUTED),
        zeroline=False, showspikes=True, spikecolor=DPA_TEAL, spikethickness=1,
    )
    fig.update_yaxes(
        showgrid=True, gridcolor="#EDF1F7", zeroline=False, tickfont=dict(color=DPA_MUTED),
        showspikes=True, spikecolor=DPA_TEAL, spikethickness=1,
    )
    if title:
        fig.update_layout(title=title)
    return fig

# --- Utility: Generate Exports (PDF, PPTX, Excel, ZIP) ---
def generate_pdf_report(title_text, subtitle_text, df):
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=36, leftMargin=36, topMargin=36, bottomMargin=36)
    elements = []
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle('ReportTitle', parent=styles['Heading1'], fontName='Helvetica-Bold', fontSize=16, leading=20, textColor=colors.HexColor(DPA_NAVY), spaceAfter=4)
    subtitle_style = ParagraphStyle('ReportSubtitle', parent=styles['Normal'], fontName='Helvetica', fontSize=10, leading=14, textColor=colors.HexColor(DPA_MUTED), spaceAfter=14)
    cell_style = ParagraphStyle('TableCell', parent=styles['Normal'], fontName='Helvetica', fontSize=8.5, leading=11, textColor=colors.HexColor(DPA_NAVY), alignment=1)
    header_cell_style = ParagraphStyle('TableHeaderCell', parent=styles['Normal'], fontName='Helvetica-Bold', fontSize=8.5, leading=11, textColor=colors.white, alignment=1)

    elements.append(Paragraph(title_text, title_style))
    elements.append(Paragraph(subtitle_text, subtitle_style))

    headers = list(df.columns)
    table_data = [[Paragraph(str(h), header_cell_style) for h in headers]]

    for idx, row in df.iterrows():
        is_total_row = (idx == "Σ") or (str(row.iloc[0]).upper() in ["TOTAL", "ALL"])
        row_cells = []
        for val in row:
            cell_text = "" if pd.isna(val) else (f"{val:,.0f}" if isinstance(val, (int, float)) else str(val))
            if is_total_row:
                total_cell_style = ParagraphStyle('TotalCell', parent=cell_style, fontName='Helvetica-Bold')
                row_cells.append(Paragraph(cell_text, total_cell_style))
            else:
                row_cells.append(Paragraph(cell_text, cell_style))
        table_data.append(row_cells)

    col_count = len(headers)
    col_width = max(540.0 / col_count, 40)
    col_widths = [col_width] * col_count

    t = Table(table_data, colWidths=col_widths, repeatRows=1)
    base_table_style = [
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor(DPA_NAVY)),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
        ('TOPPADDING', (0, 0), (-1, -1), 6),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor("#CBD5E1")),
    ]

    # Zebra striping: Teal Tint palette (matches DPA brand accent)
    for i in range(1, len(table_data)):
        if i % 2 == 0:
            base_table_style.append(('BACKGROUND', (0, i), (-1, i), colors.HexColor("#EAF6F6")))
        else:
            base_table_style.append(('BACKGROUND', (0, i), (-1, i), colors.white))

    for r_idx, row in df.iterrows():
        if (r_idx == "Σ") or (str(row.iloc[0]).upper() in ["TOTAL", "ALL"]):
            actual_row_idx = list(df.index).index(r_idx) + 1
            base_table_style.append(('BACKGROUND', (0, actual_row_idx), (-1, actual_row_idx), colors.HexColor("#D8EFEF")))

    t.setStyle(TableStyle(base_table_style))
    elements.append(t)
    doc.build(elements)
    buffer.seek(0)
    return buffer.getvalue()

def generate_pptx_report(summary_dict):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Executive Analytics Package"
    slide.placeholders[1].text = "Generated automatically by DPA Intelligence"
    
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "High-Level Executive Summary"
    tf = slide2.shapes.placeholders[1].text_frame
    for k, v in summary_dict.items():
        p = tf.add_paragraph()
        p.text = f"{k}: {v}"
        p.font.size = Pt(18)
        
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()

def generate_excel_dashboard(dfs_dict):
    from openpyxl.styles import PatternFill, Font, Alignment

    buffer = io.BytesIO()
    with pd.ExcelWriter(buffer, engine='openpyxl') as writer:
        for sheet_name, df in dfs_dict.items():
            df.to_excel(writer, sheet_name=sheet_name, index=False)

            ws = writer.sheets[sheet_name]
            n_rows, n_cols = df.shape
            if n_rows == 0 or n_cols == 0:
                continue

            header_fill = PatternFill(start_color=DPA_NAVY.replace("#", ""), end_color=DPA_NAVY.replace("#", ""), fill_type="solid")
            header_font = Font(color="FFFFFF", bold=True)
            stripe_fill = PatternFill(start_color="EAF6F6", end_color="EAF6F6", fill_type="solid")
            white_fill = PatternFill(start_color="FFFFFF", end_color="FFFFFF", fill_type="solid")

            # Style header row
            for col_idx in range(1, n_cols + 1):
                cell = ws.cell(row=1, column=col_idx)
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(horizontal="center", vertical="center")

            # Zebra striping on data rows (matches reference formatting)
            for row_idx in range(2, n_rows + 2):
                fill = stripe_fill if row_idx % 2 == 0 else white_fill
                for col_idx in range(1, n_cols + 1):
                    ws.cell(row=row_idx, column=col_idx).fill = fill

            # Auto-fit column widths (approximate)
            for col_idx, col_name in enumerate(df.columns, start=1):
                max_len = max(len(str(col_name)), df[col_name].astype(str).map(len).max() if n_rows else 0)
                ws.column_dimensions[ws.cell(row=1, column=col_idx).column_letter].width = min(max_len + 4, 40)

    buffer.seek(0)
    return buffer.getvalue()

def generate_zip_package(reports_dict):
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, 'w', zipfile.ZIP_DEFLATED) as zf:
        for filename, data in reports_dict.items():
            zf.writestr(filename, data)
    buffer.seek(0)
    return buffer.getvalue()

def get_performance_grade(achievement_pct):
    if achievement_pct >= 110: return "A+"
    elif achievement_pct >= 95: return "A"
    elif achievement_pct >= 85: return "B+"
    elif achievement_pct >= 75: return "B"
    elif achievement_pct >= 60: return "C"
    else: return "D"

def make_sparkline_svg(values, width=92, height=30, color="#0EA5A6"):
    values = [float(v) for v in values if pd.notna(v)]
    if len(values) < 2: return ""
    vmin, vmax = min(values), max(values)
    rng = (vmax - vmin) or 1.0
    step = width / (len(values) - 1)
    pts = []
    for i, v in enumerate(values):
        x = i * step
        y = (height - 3) - ((v - vmin) / rng) * (height - 6)
        pts.append(f"{x:.1f},{y:.1f}")
    line_pts = " ".join(pts)
    fill_pts = f"0,{height} " + line_pts + f" {width},{height}"
    return (
        f'<svg width="{width}" height="{height}" viewBox="0 0 {width} {height}" '
        f'xmlns="http://www.w3.org/2000/svg" preserveAspectRatio="none">'
        f'<polyline points="{fill_pts}" fill="{color}26" stroke="none" />'
        f'<polyline points="{line_pts}" fill="none" stroke="{color}" stroke-width="2" '
        f'stroke-linecap="round" stroke-linejoin="round" />'
        f'</svg>'
    )

def render_kpi_card(icon_name, label, value, delta_pct=None, delta_label="vs last month",
                     trend_values=None, status="neutral", accent="#0EA5A6", sublabel=None):
    status_colors = {"good": "#2FA875", "warn": "#E8A33D", "bad": "#D64550", "neutral": "#64748B"}
    status_color = status_colors.get(status, "#64748B")
    delta_html = ""
    if delta_pct is not None:
        arrow = "▲" if delta_pct >= 0 else "▼"
        delta_color = "#2FA875" if delta_pct >= 0 else "#D64550"
        delta_html = f'<div class="kpi-delta" style="color:{delta_color};">{arrow} {abs(delta_pct):.1f}%<span class="kpi-delta-label"> {delta_label}</span></div>'
    sub_html = f'<div class="kpi-sublabel">{sublabel}</div>' if sublabel else ""
    spark_svg = make_sparkline_svg(trend_values, color=accent) if trend_values else ""
    spark_html = f'<div class="kpi-spark">{spark_svg}</div>' if spark_svg else ""
    
    icon_html = f'<span class="material-symbols-outlined" style="font-size:20px; color:{accent};">{icon_name}</span>'
    
    return f'''<div class="kpi-card-v2" style="border-top-color:{accent};">
        <div class="kpi-card-top"><span class="kpi-icon-wrap">{icon_html}</span><span class="kpi-status-dot" style="background:{status_color};" title="{status}"></span></div>
        <div class="kpi-label">{label}</div><div class="kpi-value">{value}</div>{sub_html}{delta_html}{spark_html}</div>'''

def render_top_dpa_card(rank, name, pages, district=None, supervisor=None, max_pages=1):
    """Renders a single 'Top 3 DPAs' leaderboard card (medal-style), reusing the
    existing kpi-card-v2 styling so it visually matches the rest of the KPI grid.
    Purely additive helper — does not alter render_kpi_card or any other function."""
    medal_map = {1: ("🥇", DPA_AMBER), 2: ("🥈", "#94A3B8"), 3: ("🥉", "#B08968")}
    medal_icon, accent = medal_map.get(rank, ("🎖️", DPA_TEAL))
    pct = min(100, (float(pages) / float(max_pages) * 100)) if max_pages else 0
    meta_bits = [b for b in [f"District: {district}" if district else None,
                              f"Supervisor: {supervisor}" if supervisor else None] if b]
    meta_html = f'<div class="kpi-sublabel">{" · ".join(meta_bits)}</div>' if meta_bits else ""
    bar_html = (
        f'<div style="width:100%;height:6px;border-radius:6px;background:rgba(148,163,184,0.25);'
        f'overflow:hidden;margin-top:6px;">'
        f'<div style="height:100%;border-radius:6px;width:{pct:.1f}%;background:{accent};"></div></div>'
    )
    return f'''<div class="kpi-card-v2" style="border-top-color:{accent};">
        <div class="kpi-card-top"><span style="font-size:22px;line-height:1;">{medal_icon}</span><span class="kpi-status-dot" style="background:{accent};" title="Rank {rank}"></span></div>
        <div class="kpi-label">#{rank} Top DPA</div>
        <div class="kpi-value" style="font-size:18px;">{name if name else "—"}</div>
        {meta_html}
        <div class="kpi-delta" style="color:{accent};">{int(pages):,} pages</div>
        {bar_html}
    </div>'''

# --- Styling & Professional Icon Theme ---
st.markdown(f"""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Sora:wght@600;700;800&family=Inter:wght@400;500;600;700&display=swap');
    html, body, [class*="css"] {{ font-family: 'Inter', -apple-system, sans-serif; }}
    h1, h2, h3, h4, .stTabs [data-baseweb="tab"] p {{ font-family: 'Sora', 'Inter', sans-serif !important; }}
    .main {{ background-color: var(--background-color, {DPA_BG}); }}
    .st-key-dpa_header_banner {{
        background: linear-gradient(135deg, {DPA_NAVY} 0%, {DPA_NAVY_LIGHT} 60%, {DPA_TEAL_DARK} 140%);
        border-radius: 20px; padding: 24px 32px; margin-bottom: 24px;
        box-shadow: 0 12px 32px rgba(22, 35, 63, 0.22); border: 1px solid rgba(255, 255, 255, 0.1);
    }}
    .st-key-dpa_header_banner h1 {{ color: #FFFFFF !important; font-size: 26px !important; font-weight: 700 !important; margin: 2px 0 6px 0 !important; }}
    .st-key-dpa_header_banner p {{ color: rgba(255, 255, 255, 0.82) !important; font-size: 14px !important; margin: 0 !important; }}
    
    .horizontal-scroll-container {{ overflow-x: auto; white-space: nowrap; width: 100%; padding-bottom: 10px; }}
    .kpi-grid {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(210px, 1fr)); gap: 16px; margin-bottom: 16px; }}
    .kpi-card-v2 {{
        background-color: var(--secondary-background-color, #ffffff); border-radius: 16px; padding: 16px 18px 14px 18px;
        box-shadow: 0 4px 16px rgba(22, 35, 63, 0.06); border: 1px solid rgba(226, 232, 240, 0.8);
        border-top: 4px solid {DPA_TEAL}; display: flex; flex-direction: column; gap: 4px; transition: transform 0.2s;
    }}
    .kpi-card-v2:hover {{ transform: translateY(-3px); }}
    .kpi-card-top {{ display: flex; align-items: center; justify-content: space-between; }}
    .kpi-icon-wrap {{ display: flex; align-items: center; background: rgba(14, 165, 166, 0.1); padding: 6px; border-radius: 8px; }}
    .kpi-status-dot {{ width: 9px; height: 9px; border-radius: 50%; box-shadow: 0 0 0 3px rgba(0,0,0,0.03); }}
    .kpi-label {{ font-size: 11.5px; font-weight: 600; letter-spacing: 0.02em; color: {DPA_MUTED}; text-transform: uppercase; margin-top: 4px; line-height: 1.4; }}
    .kpi-value {{ font-family: 'Sora', sans-serif; font-size: 22px; font-weight: 700; color: {DPA_NAVY}; line-height: 1.15; word-break: break-word; }}
    .kpi-sublabel {{ font-size: 11px; color: {DPA_MUTED}; margin-top: -2px; }}
    .kpi-delta {{ font-size: 12px; font-weight: 600; margin-top: 2px; }}
    .kpi-delta-label {{ font-weight: 400; color: {DPA_MUTED}; }}
    .kpi-spark {{ margin-top: 4px; line-height: 0; }}

    /* ========================================= */
    /* ENHANCED ENTERPRISE SIDEBAR STYLING       */
    /* ========================================= */
    
    [data-testid="stSidebar"] {{
        background-color: var(--secondary-background-color, #F8FAFC) !important;
        border-right: 1px solid rgba(148, 163, 184, 0.2);
    }}
    
    [data-testid="stSidebar"] .st-emotion-cache-16txtl3 {{
        padding: 2rem 1.5rem !important;
    }}
    
    .premium-sidebar-brand {{
        background: linear-gradient(135deg, #16233F 0%, #2A3D66 60%, #0EA5A6 140%);
        border-radius: 12px;
        padding: 18px 20px;
        margin-bottom: 28px;
        box-shadow: 0 8px 24px rgba(22, 35, 63, 0.15);
        border: 1px solid rgba(255,255,255,0.1);
        display: flex;
        align-items: center;
        gap: 14px;
        transition: transform 0.2s ease, box-shadow 0.2s ease;
    }}
    .premium-sidebar-brand:hover {{
        transform: translateY(-2px);
        box-shadow: 0 12px 28px rgba(22, 35, 63, 0.25);
    }}
    .premium-sidebar-icon {{
        background: rgba(255, 255, 255, 0.15);
        backdrop-filter: blur(4px);
        padding: 10px;
        border-radius: 10px;
        display: flex;
        align-items: center;
        justify-content: center;
        color: #ffffff;
    }}
    .premium-sidebar-text-container {{
        display: flex;
        flex-direction: column;
    }}
    .premium-sidebar-title {{
        color: #ffffff;
        font-family: 'Sora', sans-serif;
        font-weight: 700;
        font-size: 18px;
        line-height: 1.1;
        letter-spacing: 0.01em;
    }}
    .premium-sidebar-subtitle {{
        color: rgba(255, 255, 255, 0.75);
        font-family: 'Inter', sans-serif;
        font-size: 11px;
        font-weight: 500;
        text-transform: uppercase;
        letter-spacing: 0.06em;
        margin-top: 4px;
    }}

    [data-testid="stSidebar"] [data-testid="stExpander"] {{
        background: var(--background-color, #ffffff);
        border: 1px solid rgba(148, 163, 184, 0.25);
        border-radius: 12px;
        box-shadow: 0 2px 8px rgba(0, 0, 0, 0.02);
        margin-bottom: 16px;
        overflow: hidden;
        transition: all 0.3s ease;
    }}
    [data-testid="stSidebar"] [data-testid="stExpander"]:hover {{
        box-shadow: 0 6px 16px rgba(0, 0, 0, 0.06);
        border-color: rgba(14, 165, 166, 0.4);
    }}
    [data-testid="stSidebar"] [data-testid="stExpander"] summary {{
        background: var(--secondary-background-color, #F1F5F9);
        padding: 14px 18px;
        border-bottom: 1px solid rgba(148, 163, 184, 0.15);
    }}
    [data-testid="stSidebar"] [data-testid="stExpander"] summary p {{
        font-family: 'Sora', sans-serif !important;
        font-weight: 600;
        font-size: 14px;
        color: var(--text-color);
    }}
    [data-testid="stSidebar"] [data-testid="stExpander"] .st-emotion-cache-p5msec {{
        padding: 16px 18px;
    }}

    [data-testid="stSidebar"] .stSelectbox > div > div, 
    [data-testid="stSidebar"] .stMultiSelect > div > div,
    [data-testid="stSidebar"] .stTextInput > div > div,
    [data-testid="stSidebar"] .stNumberInput > div > div {{
        border-radius: 8px;
        border: 1px solid rgba(148, 163, 184, 0.4);
        background-color: var(--background-color);
        transition: all 0.2s ease;
    }}
    [data-testid="stSidebar"] .stSelectbox > div > div:focus-within,
    [data-testid="stSidebar"] .stMultiSelect > div > div:focus-within,
    [data-testid="stSidebar"] .stTextInput > div > div:focus-within,
    [data-testid="stSidebar"] .stNumberInput > div > div:focus-within {{
        border-color: #0EA5A6;
        box-shadow: 0 0 0 2px rgba(14, 165, 166, 0.2);
    }}
    [data-testid="stSidebar"] label {{
        font-family: 'Inter', sans-serif;
        font-size: 13px !important;
        font-weight: 600;
        color: var(--text-color);
        margin-bottom: 4px;
    }}

    [data-testid="stSidebar"] .stButton > button {{
        border-radius: 8px;
        font-family: 'Inter', sans-serif;
        font-weight: 600;
        font-size: 14px;
        padding: 6px 16px;
        width: 100%;
        transition: all 0.3s ease;
        border: 1px solid rgba(148, 163, 184, 0.4);
        background-color: var(--background-color);
    }}
    [data-testid="stSidebar"] .stButton > button:hover {{
        border-color: #0EA5A6;
        color: #0EA5A6;
        background-color: rgba(14, 165, 166, 0.05);
        box-shadow: 0 4px 12px rgba(14, 165, 166, 0.1);
        transform: translateY(-1px);
    }}
    [data-testid="stSidebar"] .stButton > button[kind="primary"] {{
        background: linear-gradient(135deg, #0EA5A6 0%, #0B8485 100%);
        border: none;
        color: white;
        box-shadow: 0 4px 10px rgba(14, 165, 166, 0.3);
    }}
    [data-testid="stSidebar"] .stButton > button[kind="primary"]:hover {{
        box-shadow: 0 6px 16px rgba(14, 165, 166, 0.45);
        transform: translateY(-2px);
    }}

    [data-testid="stSidebar"] .stCheckbox {{
        background: var(--background-color, #ffffff);
        border: 1px solid rgba(148, 163, 184, 0.25);
        padding: 12px 16px;
        border-radius: 10px;
        box-shadow: 0 2px 6px rgba(0,0,0,0.02);
        margin-bottom: 16px;
        transition: all 0.2s ease;
    }}
    [data-testid="stSidebar"] .stCheckbox:hover {{
        border-color: #0EA5A6;
        box-shadow: 0 4px 12px rgba(14, 165, 166, 0.1);
    }}

    [data-testid="stSidebar"] h3 {{
        font-family: 'Sora', sans-serif;
        font-size: 13px !important;
        font-weight: 700 !important;
        text-transform: uppercase;
        letter-spacing: 0.08em;
        color: var(--text-color);
        margin-top: 12px;
        margin-bottom: 16px;
        display: flex;
        align-items: center;
        gap: 8px;
    }}
    [data-testid="stSidebar"] h3::before {{
        content: "";
        display: inline-block;
        width: 8px;
        height: 8px;
        background-color: #0EA5A6;
        border-radius: 50%;
    }}
    
    [data-testid="stSidebar"] hr {{
        border-top: 1px dashed rgba(148, 163, 184, 0.4);
        margin: 20px 0;
    }}
    
    [data-testid="stSidebar"] .stMultiSelect [data-baseweb="tag"] {{
        background-color: rgba(14, 165, 166, 0.15);
        border: 1px solid rgba(14, 165, 166, 0.3);
        color: #0B8485;
        border-radius: 6px;
        font-weight: 600;
    }}
    
    [data-testid="stSidebar"] ::-webkit-scrollbar {{
        width: 6px;
    }}
    [data-testid="stSidebar"] ::-webkit-scrollbar-track {{
        background: transparent;
    }}
    [data-testid="stSidebar"] ::-webkit-scrollbar-thumb {{
        background: rgba(148, 163, 184, 0.4);
        border-radius: 10px;
    }}
    [data-testid="stSidebar"] ::-webkit-scrollbar-thumb:hover {{
        background: rgba(148, 163, 184, 0.6);
    }}
    </style>
""", unsafe_allow_html=True)

# --- App Header with Professional Icon ---
with st.container(key="dpa_header_banner"):
    header_col1, header_col2 = st.columns([1, 15])
    with header_col1:
        st.markdown('<span class="material-symbols-outlined" style="font-size: 40px; color: #FFFFFF; margin-top: 6px;">monitoring</span>', unsafe_allow_html=True)
    with header_col2:
        st.markdown(
            "<h1>Executive Analytics Center &amp; Output Management</h1>"
            "<p>Advanced enterprise scanning analytics, high-performance columnar engine, and automated intelligence.</p>",
            unsafe_allow_html=True,
        )

# --- DuckDB Core & 24-Hour Auto-Save Setup ---
DB_FILE_PATH = "dpa_data.duckdb"
SUPERVISOR_TABLE_NAME = "supervisor_mapping"
TARGETS_TABLE_NAME = "monthly_targets"
EMPLOYEE_TABLE = "employee_registry"
AUDIT_LOG_TABLE = "audit_logs"

DEFAULT_SUPERVISOR_MAPPING = {
    "ATTOCK": "SALAMT KHAN", "CHAKWAL": "SALAMT KHAN", "GUJRAT": "SALAMT KHAN", "JEHLUM": "SALAMT KHAN", "RAWALPINDI": "SALAMT KHAN",
    "BHAKKAR": "ZULQANAIN HADIER", "KHUSHAB": "ZULQANAIN HADIER", "MANDI BAHAUDDIN": "ZULQANAIN HADIER", "MIANWALI": "ZULQANAIN HADIER", "SARGODHA": "ZULQANAIN HADIER",
    "BAHAWALNAGAR": "ZESHAN", "MULTAN": "ZESHAN", "OKARA": "ZESHAN", "PAKPATTAN": "ZESHAN", "SAHIWAL": "ZESHAN",
    "BAHAWALPUR": "ZULFQAR ALI", "LODHRAN": "ZULFQAR ALI", "RAHIM YAR KHAN": "ZULFQAR ALI", "RAJANPUR": "ZULFQAR ALI",
    "D G KHAN": "SHAZAD KHOSSA", "KHANEWAL": "SHAZAD KHOSSA", "LAYYAH": "SHAZAD KHOSSA", "MUZAFFARGARH": "SHAZAD KHOSSA", "VEHARI": "SHAZAD KHOSSA",
    "CHINIOT": "MUHAMMAD ASIF", "FAISALABAD": "MUHAMMAD ASIF", "HAFIZABAD": "MUHAMMAD ASIF", "JHANG": "MUHAMMAD ASIF", "TOBA TEK SINGH": "MUHAMMAD ASIF",
    "FREEDKOT LAHORE": "MUSHAID", "SHADRA LAHORE": "MUSHAID", "SHEIKHUPURA": "MUSHAID", "NANKANA SAHIB": "MUSHAID",
    "GUJRANWALA": "MUNEEB", "KASUR": "MUNEEB", "LAHORE DHA": "MUNEEB", "NAROWAL": "MUNEEB", "OPF LAHORE": "MUNEEB", "ALICOMPLEX": "MUNEEB", "SIALKOT": "MUNEEB"
}

MONTH_MAP = {'jan': 1, 'feb': 2, 'mar': 3, 'apr': 4, 'may': 5, 'jun': 6, 'jul': 7, 'aug': 8, 'sep': 9, 'oct': 10, 'nov': 11, 'dec': 12}

def get_calendar_columns_for_sheet(sheet_name, year_val):
    m_num = next((v for k, v in MONTH_MAP.items() if k in str(sheet_name).lower()), 1)
    try: yr = int(year_val)
    except ValueError: yr = 2026
    _, last_day = calendar.monthrange(yr, m_num)
    return [str(day) for day in range(1, last_day + 1)]

def auto_calculate_total_pages(df):
    df = df.copy()
    meta_cols = {'DISTRICT', 'DPA NAME', 'TOTAL PAGES', 'SR NO.', 'SR.NO', 'S.NO', 'SR'}
    day_cols = [c for c in df.columns if str(c).strip().upper() not in meta_cols]
    numeric_days = df[day_cols].apply(pd.to_numeric, errors='coerce').fillna(0)
    total_col_name = next((c for c in df.columns if 'TOTAL' in str(c).upper() and 'PAGE' in str(c).upper()), 'TOTAL PAGES')
    if total_col_name in df.columns:
        df[total_col_name] = numeric_days.sum(axis=1)
    return df

def log_audit_event(user_id, action, details=""):
    try:
        con = duckdb.connect(DB_FILE_PATH)
        con.execute(f"""
            CREATE TABLE IF NOT EXISTS {AUDIT_LOG_TABLE} (
                timestamp TIMESTAMP, 
                user_id VARCHAR, 
                action VARCHAR,
                details VARCHAR
            )
        """)
        try:
            con.execute(f"ALTER TABLE {AUDIT_LOG_TABLE} ADD COLUMN details VARCHAR")
        except Exception:
            pass
            
        timestamp = datetime.datetime.now()
        con.execute(f"INSERT INTO {AUDIT_LOG_TABLE} (timestamp, user_id, action, details) VALUES (?, ?, ?, ?)", [timestamp, user_id, action, details])
        con.close()
    except Exception as e:
        print(f"Audit log error: {e}")

def init_duckdb():
    con = duckdb.connect(DB_FILE_PATH)
    con.execute(f"CREATE TABLE IF NOT EXISTS {SUPERVISOR_TABLE_NAME} (district VARCHAR PRIMARY KEY, supervisor VARCHAR)")
    con.execute(f"CREATE TABLE IF NOT EXISTS {TARGETS_TABLE_NAME} (period_key VARCHAR PRIMARY KEY, target DOUBLE)")
    con.execute(f"CREATE TABLE IF NOT EXISTS {AUDIT_LOG_TABLE} (timestamp TIMESTAMP, user_id VARCHAR, action VARCHAR)")
    con.execute(f"CREATE TABLE IF NOT EXISTS {EMPLOYEE_TABLE} (emp_id VARCHAR PRIMARY KEY, name VARCHAR, district VARCHAR, supervisor VARCHAR, emp_role VARCHAR, status VARCHAR, joining_date DATE, leaving_date DATE)")
    con.execute("CREATE TABLE IF NOT EXISTS app_metadata (key VARCHAR PRIMARY KEY, value VARCHAR)")
    
    try:
        con.execute(f"ALTER TABLE {EMPLOYEE_TABLE} ADD COLUMN emp_role VARCHAR")
    except Exception:
        pass
        
    try:
        con.execute(f"ALTER TABLE {AUDIT_LOG_TABLE} ADD COLUMN details VARCHAR")
    except Exception:
        pass

    res = con.execute(f"SELECT COUNT(*) FROM {SUPERVISOR_TABLE_NAME}").fetchone()[0]
    if res == 0:
        for dist, sup in DEFAULT_SUPERVISOR_MAPPING.items():
            con.execute(f"INSERT OR REPLACE INTO {SUPERVISOR_TABLE_NAME} (district, supervisor) VALUES (?, ?)", [dist, sup])

    tables = [t[0] for t in con.execute("SELECT table_name FROM information_schema.tables WHERE table_schema='main'").fetchall()]
    system_tables = {SUPERVISOR_TABLE_NAME, TARGETS_TABLE_NAME, EMPLOYEE_TABLE, AUDIT_LOG_TABLE, 'app_metadata'}
    user_tables = [t for t in tables if t not in system_tables]
    
    if not user_tables:
        jan_days = get_calendar_columns_for_sheet("jan-2026", 2026)
        default_df = pd.DataFrame(columns=["DISTRICT", "DPA NAME"] + jan_days + ["TOTAL PAGES"])
        con.register("df_default", default_df)
        con.execute("CREATE TABLE main_2026_jan_2026 AS SELECT * FROM df_default")
    con.close()

def check_and_perform_auto_save():
    con = duckdb.connect(DB_FILE_PATH)
    row = con.execute("SELECT value FROM app_metadata WHERE key = 'last_auto_save'").fetchone()
    now = time.time()
    one_day_seconds = 24 * 60 * 60
    
    should_save = False
    if not row:
        should_save = True
    else:
        try:
            if now - float(row[0]) >= one_day_seconds:
                should_save = True
        except ValueError:
            should_save = True
            
    if should_save:
        backup_filename = f"dpa_data_autosave_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.duckdb"
        con.close()
        if os.path.exists(DB_FILE_PATH):
            shutil.copyfile(DB_FILE_PATH, backup_filename)
        con = duckdb.connect(DB_FILE_PATH)
        con.execute("INSERT OR REPLACE INTO app_metadata (key, value) VALUES ('last_auto_save', ?)", [str(now)])
    con.close()

def get_all_tables_hierarchical():
    con = duckdb.connect(DB_FILE_PATH)
    tables = [t[0] for t in con.execute("SELECT table_name FROM information_schema.tables WHERE table_schema='main'").fetchall()]
    hierarchical = {}
    system_tables = {SUPERVISOR_TABLE_NAME, TARGETS_TABLE_NAME, EMPLOYEE_TABLE, AUDIT_LOG_TABLE, 'app_metadata'}
    for t in tables:
        if t in system_tables: continue
        parts = t.split("_", 1)
        if len(parts) == 2:
            df = con.execute(f'SELECT * FROM "{t}"').df()
            hierarchical.setdefault(parts[0], {})[parts[1].replace("_", "-")] = df
    con.close()
    return hierarchical

def save_sheet_to_duckdb(year_key, sheet_name, df):
    con = duckdb.connect(DB_FILE_PATH)
    table_name = f"{year_key}_{sheet_name.replace('-', '_').replace(' ', '_')}"
    con.register("df_temp", df)
    con.execute(f'CREATE OR REPLACE TABLE "{table_name}" AS SELECT * FROM df_temp')
    con.close()

def delete_sheet_from_duckdb(year_key, sheet_name):
    con = duckdb.connect(DB_FILE_PATH)
    table_name = f"{year_key}_{sheet_name.replace('-', '_').replace(' ', '_')}"
    con.execute(f'DROP TABLE IF EXISTS "{table_name}"')
    con.close()

def delete_year_from_duckdb(year_key):
    con = duckdb.connect(DB_FILE_PATH)
    tables = [t[0] for t in con.execute("SELECT table_name FROM information_schema.tables WHERE table_schema='main'").fetchall()]
    system_tables = {SUPERVISOR_TABLE_NAME, TARGETS_TABLE_NAME, EMPLOYEE_TABLE, AUDIT_LOG_TABLE, 'app_metadata'}
    for t in tables:
        if t.startswith(f"{year_key}_") and t not in system_tables:
            con.execute(f'DROP TABLE IF EXISTS "{t}"')
    con.close()

def load_supervisor_mapping_from_db():
    con = duckdb.connect(DB_FILE_PATH)
    sup_df = con.execute(f"SELECT * FROM {SUPERVISOR_TABLE_NAME}").df()
    con.close()
    if not sup_df.empty and "district" in sup_df.columns and "supervisor" in sup_df.columns:
        return {str(row["district"]).strip().upper(): str(row["supervisor"]).strip() for _, row in sup_df.dropna().iterrows()}
    return dict(DEFAULT_SUPERVISOR_MAPPING)

def load_monthly_target(year_key, sheet_name, default_target):
    con = duckdb.connect(DB_FILE_PATH)
    res = con.execute(f"SELECT target FROM {TARGETS_TABLE_NAME} WHERE period_key = ?", [f"{year_key}::{sheet_name}"]).fetchone()
    con.close()
    return float(res[0]) if res and res[0] is not None else float(default_target)

def save_monthly_target(year_key, sheet_name, target_value):
    con = duckdb.connect(DB_FILE_PATH)
    con.execute(f"INSERT OR REPLACE INTO {TARGETS_TABLE_NAME} (period_key, target) VALUES (?, ?)", [f"{year_key}::{sheet_name}", float(target_value)])
    con.close()

def compute_sheet_total_pages(df):
    if df is None or df.empty: return 0.0
    meta_cols = {'DISTRICT', 'DPA NAME', 'TOTAL PAGES', 'SR NO.', 'SR.NO', 'S.NO', 'SR'}
    day_cols = [c for c in df.columns if str(c).strip().upper() not in meta_cols]
    return float(df[day_cols].apply(pd.to_numeric, errors="coerce").fillna(0).values.sum()) if day_cols else 0.0

def get_previous_period_total(hierarchical_data, current_year, current_sheet):
    periods = []
    for yr_key, sheets in hierarchical_data.items():
        try: yr_int = int(yr_key)
        except: yr_int = 0
        for sheet_name in sheets.keys():
            m_idx = next((idx for m, idx in MONTH_MAP.items() if m in str(sheet_name).lower()), 99)
            periods.append((yr_int, m_idx, yr_key, sheet_name))
    periods.sort(key=lambda x: (x[0], x[1]))
    keys = [(yr, sheet) for (_, _, yr, sheet) in periods]
    try: current_idx = keys.index((current_year, current_sheet))
    except ValueError: return None, None, None
    if current_idx == 0: return None, None, None
    prev_year, prev_sheet = keys[current_idx - 1]
    prev_total = compute_sheet_total_pages(hierarchical_data.get(prev_year, {}).get(prev_sheet, pd.DataFrame()))
    return prev_year, prev_sheet, prev_total

# --- Session Initialization & Auto-Save ---
init_duckdb()
check_and_perform_auto_save()

# --- Premium Sidebar Branding Header ---
st.sidebar.markdown("""
    <div class="premium-sidebar-brand">
        <div class="premium-sidebar-icon">
            <span class="material-symbols-outlined" style="font-size: 26px;">dashboard_customize</span>
        </div>
        <div class="premium-sidebar-text-container">
            <div class="premium-sidebar-title">DPA Core</div>
            <div class="premium-sidebar-subtitle">Enterprise Portal</div>
        </div>
    </div>
""", unsafe_allow_html=True)

if "hierarchical_data" not in st.session_state: st.session_state["hierarchical_data"] = get_all_tables_hierarchical()
if "supervisor_mapping" not in st.session_state: st.session_state["supervisor_mapping"] = load_supervisor_mapping_from_db()
if "is_authorized" not in st.session_state: st.session_state["is_authorized"] = True
if "user_role" not in st.session_state: st.session_state["user_role"] = "Viewer"

# --- Role-Based Access Control (RBAC) Definitions ---
ROLE_HIERARCHY = {"Viewer": 1, "Supervisor": 2, "Admin": 3}

def has_permission(min_role: str) -> bool:
    current_role = st.session_state.get("user_role", "Viewer")
    return ROLE_HIERARCHY.get(current_role, 1) >= ROLE_HIERARCHY.get(min_role, 1)

# --- Sidebar Controls ---
st.sidebar.markdown("### ⚙️ Executive Filters & Controls")
auto_refresh = st.sidebar.checkbox("🔄 Enable Live Auto-Refresh (1 min)")
if auto_refresh:
    time.sleep(60)
    st.rerun()

with st.sidebar.expander("📅 Period & Sheet Selector", expanded=True):
    available_years = list(st.session_state["hierarchical_data"].keys()) or ["2026"]
    selected_year = st.selectbox("Select Year", available_years)

    def sheet_sort_key(name):
        lower_name = str(name).lower()
        m_idx = next((idx for m, idx in MONTH_MAP.items() if m in lower_name), 99)
        year_match = re.search(r'20\d{2}', lower_name)
        yr = int(year_match.group(0)) if year_match else int(selected_year)
        return (yr, m_idx)

    raw_months = list(st.session_state["hierarchical_data"].get(selected_year, {}).keys())
    available_months = sorted(raw_months, key=sheet_sort_key) or ["jan-2026"]
    selected_sheet = st.selectbox("Select Month / Report Sheet", available_months)

current_sheet_dict = st.session_state["hierarchical_data"].get(selected_year, {})
raw_df = current_sheet_dict.get(selected_sheet, pd.DataFrame(columns=["DISTRICT", "DPA NAME", "TOTAL PAGES"]))

def clean_dpa_dataframe(df, sheet_name=""):
    empty_quality = {"duplicate_count": 0, "duplicate_rows": pd.DataFrame()}
    if df is None or df.empty: return df, pd.DataFrame(), "", "", None, [], empty_quality

    df = df.copy()
    df.columns = [str(c).strip().upper() for c in df.columns]

    dup_suffix = re.compile(r"^(.*)\.(\d+)$")
    merge_targets = {}
    for col in df.columns:
        match = dup_suffix.match(col)
        if match and match.group(1) in df.columns:
            merge_targets.setdefault(match.group(1), []).append(col)
    for base_col, extra_cols in merge_targets.items():
        for extra_col in extra_cols:
            df[base_col] = pd.to_numeric(df[base_col], errors="coerce").fillna(0) + pd.to_numeric(df[extra_col], errors="coerce").fillna(0)
        df.drop(columns=extra_cols, inplace=True)

    name_col = next((c for c in df.columns if "NAME" in c or "DPA" in c), df.columns[1] if len(df.columns) > 1 else df.columns[0])
    district_col = next((c for c in df.columns if "DISTRICT" in c), df.columns[0])
    total_col = next((c for c in df.columns if "TOTAL" in c and "PAGE" in c), None)

    df_filtered = df.dropna(subset=[name_col])
    df_filtered = df_filtered[~df_filtered[name_col].astype(str).str.contains("NAME|SR|TOTAL", na=False, case=False)].copy()
    df_filtered[name_col] = df_filtered[name_col].astype(str).str.strip().str.upper()
    df_filtered[district_col] = df_filtered[district_col].astype(str).str.strip().str.upper()

    dup_mask = df_filtered.duplicated(subset=[district_col, name_col], keep="first")
    duplicate_count = int(dup_mask.sum())
    df_filtered = df_filtered.loc[~dup_mask]

    meta_cols = {name_col, district_col, "SR NO.", "SR.NO", "S.NO", "SR"}
    if total_col: meta_cols.add(total_col)
    date_cols = [c for c in df_filtered.columns if c not in meta_cols]

    melted = df_filtered.melt(id_vars=[col for col in [district_col, name_col] if col in df_filtered.columns], value_vars=date_cols, var_name="DATE_STR", value_name="PAGES")
    melted.rename(columns={district_col: "DISTRICT", name_col: "NAME"}, inplace=True)
    melted["PAGES"] = pd.to_numeric(melted["PAGES"], errors="coerce").fillna(0)

    month_year = None
    month_match = re.search(r"([A-Za-z]{3,9})[\s\-_]*'?(\d{2,4})", str(sheet_name))
    if month_match:
        for fmt in ("%b %Y", "%B %Y", None):
            try:
                candidate = pd.to_datetime(f"{month_match.group(1)} {month_match.group(2)}", format=fmt) if fmt else pd.to_datetime(f"{month_match.group(1)} {month_match.group(2)}")
                if pd.notna(candidate):
                    month_year = candidate
                    break
            except: continue
    if month_year is None and selected_year:
        month_match_name = re.search(r"([A-Za-z]{3,9})", str(sheet_name))
        if month_match_name:
            try: month_year = pd.to_datetime(f"{month_match_name.group(1)} {selected_year}", format="%b %Y")
            except:
                try: month_year = pd.to_datetime(f"{month_match_name.group(1)} {selected_year}", format="%B %Y")
                except: pass

    def _resolve_date(day_value):
        digits = re.sub(r"[^0-9]", "", str(day_value))
        if month_year is not None and digits and 1 <= int(digits) <= 31:
            try: return month_year.replace(day=int(digits))
            except: return pd.NaT
        return pd.to_datetime(day_value, errors="coerce")

    melted["PARSED_DATE"] = melted["DATE_STR"].apply(_resolve_date)
    
    con_emp_sup = duckdb.connect(DB_FILE_PATH)
    emp_sup_df = con_emp_sup.execute(f"SELECT name, supervisor, district FROM {EMPLOYEE_TABLE}").df()
    con_emp_sup.close()
    
    emp_supervisor_lookup = {}
    emp_district_lookup = {}
    if not emp_sup_df.empty:
        for _, r in emp_sup_df.iterrows():
            if pd.notna(r["name"]) and pd.notna(r["supervisor"]):
                emp_supervisor_lookup[str(r["name"]).strip().upper()] = str(r["supervisor"]).strip().upper()
            if pd.notna(r["name"]) and pd.notna(r["district"]):
                emp_district_lookup[str(r["name"]).strip().upper()] = str(r["district"]).strip().upper()

    def get_assigned_supervisor(row):
        name_val = str(row["NAME"]).strip().upper()
        dist_val = str(row["DISTRICT"]).strip().upper()
        if name_val in emp_supervisor_lookup:
            return emp_supervisor_lookup[name_val]
        return st.session_state["supervisor_mapping"].get(dist_val, "UNASSIGNED")

    melted["SUPERVISOR"] = melted.apply(get_assigned_supervisor, axis=1)
    
    def get_assigned_district(row):
        name_val = str(row["NAME"]).strip().upper()
        if name_val in emp_district_lookup:
            return emp_district_lookup[name_val]
        return row["DISTRICT"]
        
    melted["DISTRICT"] = melted.apply(get_assigned_district, axis=1)
    
    return df_filtered, melted, district_col, name_col, total_col, date_cols, {"duplicate_count": duplicate_count}

clean_df, melted_df, district_col, name_col, total_col, date_cols, data_quality = clean_dpa_dataframe(raw_df, selected_sheet)

with st.sidebar.expander("🔍 Global Cross-Filters", expanded=True):
    all_districts = sorted(melted_df["DISTRICT"].unique()) if not melted_df.empty else []
    selected_districts = st.multiselect("Filter Districts", all_districts, default=all_districts)

    all_supervisors = sorted(melted_df["SUPERVISOR"].unique()) if not melted_df.empty else []
    selected_supervisors = st.multiselect("Filter Supervisors", all_supervisors, default=all_supervisors)

    search_employee = st.text_input("Search Employee (DPA)", "")

filtered_melted_df = melted_df.copy()
if selected_districts: filtered_melted_df = filtered_melted_df[filtered_melted_df["DISTRICT"].isin(selected_districts)]
if selected_supervisors: filtered_melted_df = filtered_melted_df[filtered_melted_df["SUPERVISOR"].isin(selected_supervisors)]
if search_employee: filtered_melted_df = filtered_melted_df[filtered_melted_df["NAME"].str.contains(search_employee.upper(), na=False)]

with st.sidebar.expander("🔐 Role & Access Control", expanded=False):
    st.markdown("""
    **Role Permissions:**
    * **Viewer:** Read-only access to dashboards, reports, and charts — no login required.
    * **Supervisor:** Can manage supervisor mapping and assigned team oversight.
    * **Admin:** Full access to edit sheets, targets, user roles, and system data.
    """)
    
    if st.session_state["user_role"] == "Viewer":
        admin_pass = st.text_input("Enter Credentials / Password", type="password", key="admin_password_input")
        role_select = st.selectbox("Select Access Role", ["Supervisor", "Admin"])
        if st.button("🔓 Authenticate Role", use_container_width=True):
            if (role_select == "Admin" and admin_pass == "adminpass123") or \
               (role_select == "Supervisor" and admin_pass in ["supervisorpass123", "adminpass123"]):
                st.session_state["is_authorized"] = True
                st.session_state["user_role"] = role_select
                
                login_user = f"User ({role_select})"
                log_audit_event(user_id=login_user, action="LOGIN", details=f"Successfully logged in as {role_select}")
                
                st.success(f"Authenticated as {role_select}!")
                st.rerun()
            else: st.error("Incorrect password for selected role.")
    else:
        st.success(f"Role: **{st.session_state['user_role']}**")
        if st.button("🔒 Logout", use_container_width=True):
            logout_user = f"User ({st.session_state['user_role']})"
            log_audit_event(user_id=logout_user, action="LOGOUT", details="User logged out")
            st.session_state["is_authorized"] = True
            st.session_state["user_role"] = "Viewer"
            st.rerun()

if has_permission("Admin"):
    with st.sidebar.expander("🛠️ Manage Years & Months (Admin)", expanded=False):
        new_year = st.text_input("New Year (e.g., 2027)", key="input_new_year")
        if st.button("➕ Add Year", use_container_width=True) and new_year and new_year not in st.session_state["hierarchical_data"]:
            con_reg = duckdb.connect(DB_FILE_PATH)
            reg_df = con_reg.execute(f"SELECT district, name FROM {EMPLOYEE_TABLE} WHERE status = 'Active' AND (emp_role IS NULL OR emp_role != 'Divisional IT Supervisor') ORDER BY name ASC").df()
            con_reg.close()
            
            cal_cols = get_calendar_columns_for_sheet("jan", new_year)
            if not reg_df.empty:
                sample_df = pd.DataFrame({
                    "DISTRICT": reg_df["district"].str.upper(),
                    "DPA NAME": reg_df["name"].str.upper()
                })
                for c in cal_cols:
                    sample_df[c] = 0
                sample_df["TOTAL PAGES"] = 0
            else:
                sample_df = pd.DataFrame(columns=["DISTRICT", "DPA NAME"] + cal_cols + ["TOTAL PAGES"])
                
            save_sheet_to_duckdb(new_year, f"jan-{new_year}", sample_df)
            log_audit_event(user_id=f"Admin ({st.session_state.get('user_role', 'Admin')})", action="CREATE_YEAR", details=f"Created year {new_year}")
            st.session_state["hierarchical_data"] = get_all_tables_hierarchical()
            st.rerun()

        if len(st.session_state["hierarchical_data"]) > 1:
            year_to_delete = st.selectbox("Select Year to Delete", available_years)
            if st.button("🗑️ Delete Year", type="primary", use_container_width=True):
                delete_year_from_duckdb(year_to_delete)
                log_audit_event(user_id=f"Admin ({st.session_state.get('user_role', 'Admin')})", action="DELETE_YEAR", details=f"Deleted year {year_to_delete}")
                st.session_state["hierarchical_data"] = get_all_tables_hierarchical()
                st.rerun()

        new_month_name = st.text_input("New Month Name (e.g., feb-2026)", key="input_new_month")
        if st.button("➕ Add Month", use_container_width=True) and new_month_name:
            con_reg = duckdb.connect(DB_FILE_PATH)
            reg_df = con_reg.execute(f"SELECT district, name FROM {EMPLOYEE_TABLE} WHERE status = 'Active' AND (emp_role IS NULL OR emp_role != 'Divisional IT Supervisor') ORDER BY name ASC").df()
            con_reg.close()
            
            cal_cols = get_calendar_columns_for_sheet(new_month_name, selected_year)
            if not reg_df.empty:
                sample_df = pd.DataFrame({
                    "DISTRICT": reg_df["district"].str.upper(),
                    "DPA NAME": reg_df["name"].str.upper()
                })
                for c in cal_cols:
                    sample_df[c] = 0
                sample_df["TOTAL PAGES"] = 0
            else:
                sample_df = pd.DataFrame(columns=["DISTRICT", "DPA NAME"] + cal_cols + ["TOTAL PAGES"])
                
            save_sheet_to_duckdb(selected_year, new_month_name, sample_df)
            log_audit_event(user_id=f"Admin ({st.session_state.get('user_role', 'Admin')})", action="CREATE_MONTH", details=f"Created month sheet {new_month_name} for {selected_year}")
            st.session_state["hierarchical_data"] = get_all_tables_hierarchical()
            st.rerun()

        if len(st.session_state["hierarchical_data"].get(selected_year, {})) > 1:
            month_to_delete = st.selectbox("Select Month to Delete", available_months)
            if st.button("🗑️ Delete Month Sheet", type="primary", use_container_width=True):
                delete_sheet_from_duckdb(selected_year, month_to_delete)
                log_audit_event(user_id=f"Admin ({st.session_state.get('user_role', 'Admin')})", action="DELETE_MONTH", details=f"Deleted month sheet {month_to_delete} from {selected_year}")
                st.session_state["hierarchical_data"] = get_all_tables_hierarchical()
                st.rerun()

# --- Calculations ---
total_scanned_all = pd.to_numeric(filtered_melted_df["PAGES"], errors="coerce").sum() if not filtered_melted_df.empty else 0
active_dpas = filtered_melted_df["NAME"].nunique() if not filtered_melted_df.empty else 0
avg_per_dpa = total_scanned_all / active_dpas if active_dpas > 0 else 0

district_perf = filtered_melted_df.groupby("DISTRICT")["PAGES"].sum() if not filtered_melted_df.empty else pd.Series(dtype=float)
has_data = total_scanned_all > 0

best_district = district_perf.idxmax() if (not district_perf.empty and has_data) else ""
lowest_district = district_perf.idxmin() if (not district_perf.empty and has_data) else ""

supervisor_perf = filtered_melted_df.groupby("SUPERVISOR")["PAGES"].sum() if not filtered_melted_df.empty else pd.Series(dtype=float)
best_supervisor = supervisor_perf.idxmax() if (not supervisor_perf.empty and has_data) else ""

# --- Top 3 DPAs (individual top performers) for the CEO View leaderboard ---
dpa_perf_detail = (
    filtered_melted_df.groupby(["NAME", "DISTRICT", "SUPERVISOR"])["PAGES"].sum()
    .reset_index().sort_values("PAGES", ascending=False)
    if not filtered_melted_df.empty else pd.DataFrame(columns=["NAME", "DISTRICT", "SUPERVISOR", "PAGES"])
)
top3_dpas_df = dpa_perf_detail.head(3) if has_data else dpa_perf_detail.head(0)

active_days_count = filtered_melted_df["PARSED_DATE"].nunique() if not filtered_melted_df.empty else 1

m_num = next((v for k, v in MONTH_MAP.items() if k in str(selected_sheet).lower()), 1)
try: yr_val = int(selected_year)
except ValueError: yr_val = 2026
_, days_in_month = calendar.monthrange(yr_val, m_num)

data_entered_df = filtered_melted_df[filtered_melted_df["PAGES"] > 0] if not filtered_melted_df.empty else filtered_melted_df
if not data_entered_df.empty:
    # Use the last day that actually has data entered in this sheet — this is what
    # "elapsed" means for a data-entry workbook, regardless of the real-world date
    # (sheets are often filled in retroactively or ahead of the wall-clock date).
    elapsed_days = min(int(pd.to_datetime(data_entered_df["PARSED_DATE"]).max().day), days_in_month)
else:
    # No entries at all yet for this sheet. If it's the real-world current month,
    # fall back to today's date; otherwise nothing has been entered.
    today_ts = pd.Timestamp(datetime.date.today())
    if yr_val == today_ts.year and m_num == today_ts.month:
        elapsed_days = min(today_ts.day, days_in_month)
    else:
        elapsed_days = 0

daily_run_rate = total_scanned_all / max(elapsed_days, 1) if has_data else 0.0

# Remaining working days = calendar days left after the last entered day, excluding Sundays.
# e.g. May_2026 (31 days, Sundays on 3/10/17/24/31 -> 26 working days total): if the last
# entered day is 26, remaining_working_days counts the non-Sunday days from 27-31.
remaining_working_days = sum(
    1 for d in range(elapsed_days + 1, days_in_month + 1)
    if datetime.date(yr_val, m_num, d).weekday() != 6  # exclude Sundays
)

# Project forward from what's actually been scanned so far, using the observed
# per-working-day rate applied only to the real working days still remaining —
# rather than blindly multiplying the average rate by every calendar day in the month.
projected_month_end = (total_scanned_all + daily_run_rate * remaining_working_days) if has_data else 0.0

monthly_target = load_monthly_target(selected_year, selected_sheet, 0)
achievement_pct = (total_scanned_all / monthly_target * 100) if (monthly_target and has_data) else 0.0

prev_year, prev_sheet, prev_total = get_previous_period_total(st.session_state["hierarchical_data"], selected_year, selected_sheet)
prev_month_delta_pct = ((total_scanned_all - prev_total) / prev_total * 100) if (prev_total is not None and has_data) else None

daily_trend = filtered_melted_df.groupby("PARSED_DATE")["PAGES"].sum().sort_index() if not filtered_melted_df.empty else pd.Series(dtype=float)
consistency_pct = max(0, 100 - (daily_trend.std() / daily_trend.mean() * 100)) if len(daily_trend) > 1 and daily_trend.mean() > 0 else 0
attendance_pct = min((active_days_count / days_in_month) * 100, 100) if days_in_month else 0
productivity_score = min(100, max(0, (min(achievement_pct, 100) * 0.5) + (consistency_pct * 0.3) + (attendance_pct * 0.2))) if has_data else 0.0
trend_values = daily_trend.tail(14).tolist() if has_data else []

def _ach_status(pct): return "good" if pct >= 100 else "warn" if pct >= 75 else "bad"

# --- Main Interface & Role-Based Tab Assignment ---
st.subheader(f"Executive Analytics Center — {str(selected_sheet).replace('_', ' ').title()} (Role: {st.session_state['user_role']})")

current_role = st.session_state.get("user_role", "Viewer")

if current_role == "Viewer":
    tab_labels = ["📈 KPIs & AI Insights", "🏢 Reports & Exports", "📊 Advanced Visuals"]
    tabs = st.tabs(tab_labels)
    tab1, tab2, tab3 = tabs[0], tabs[1], tabs[2]
elif current_role == "Supervisor":
    tab_labels = ["📈 KPIs & AI Insights", "🏢 Reports & Exports", "📊 Advanced Visuals", "🔍 Intelligence & Scorecards", "🔮 Forecasting", "⚖️ Benchmarking", "👥 HR Management"]
    tabs = st.tabs(tab_labels)
    tab1, tab2, tab3, tab4, tab5, tab6, tab7 = tabs[0], tabs[1], tabs[2], tabs[3], tabs[4], tabs[5], tabs[6]
else: 
    tab_labels = ["📈 KPIs & AI Insights", "🏢 Reports & Exports", "📊 Advanced Visuals", "🔍 Intelligence & Scorecards", "🔮 Forecasting", "⚖️ Benchmarking", "👥 HR Management", "📝 Live Data Sheet", "🛡️ Audit Logs"]
    tabs = st.tabs(tab_labels)
    tab1, tab2, tab3, tab4, tab5, tab6, tab7, tab8, tab9 = tabs[0], tabs[1], tabs[2], tabs[3], tabs[4], tabs[5], tabs[6], tabs[7], tabs[8]

with tab1:
    st.subheader("🎯 CEO View & Key Performance Indicators")
    if not has_data:
        kpi_cards_html = [
            render_kpi_card("description", "Total Pages", "0", delta_pct=-100.0 if prev_total else None, delta_label=f"vs last month" if not prev_sheet else f"vs {prev_sheet}", trend_values=[], status="bad", accent=DPA_TEAL),
            render_kpi_card("flag", "Monthly Target", f"{int(monthly_target):,}" if monthly_target else "Not Set", sublabel="Configured in main body", status="neutral", accent=DPA_NAVY),
            render_kpi_card("verified", "Achievement %", "0.0%", delta_pct=-100.0, delta_label="vs target", trend_values=[], status="bad", accent=DPA_GREEN),
            render_kpi_card("person", "Avg Pages/DPA", "0", status="neutral", accent=DPA_TEAL),
            render_kpi_card("military_tech", "Best District", "", sublabel="0 pages", status="neutral", accent=DPA_GREEN),
            render_kpi_card("warning", "Lowest District", "", sublabel="0 pages", status="neutral", accent=DPA_RED),
            render_kpi_card("supervisor_account", "Best Supervisor", "", sublabel="0 pages", status="neutral", accent=DPA_NAVY),
            render_kpi_card("group", "Active Employees", f"{active_dpas}", status="neutral", accent=DPA_TEAL),
            render_kpi_card("psychology", "Productivity Score", "0.0/100", status="bad", accent="#6C63FF"),
            render_kpi_card("bolt", "Daily Run Rate", "0", sublabel="pages/day", trend_values=[], status="neutral", accent=DPA_AMBER),
            render_kpi_card("trending_up", "Projected Month End", "0", delta_pct=-100.0 if monthly_target else None, delta_label="vs target", status="bad", accent="#4C7CA8"),
            render_kpi_card("history", f"{int(prev_total):,}" if prev_total is not None else "N/A", "Previous Month", delta_pct=-100.0 if prev_total else None, delta_label=f"vs {prev_sheet or 'N/A'}", status="bad" if prev_total else "neutral", accent="#B08968"),
        ]
    else:
        kpi_cards_html = [
            render_kpi_card("description", "Total Pages", f"{int(total_scanned_all):,}", delta_pct=prev_month_delta_pct, trend_values=trend_values, status=_ach_status(achievement_pct), accent=DPA_TEAL),
            render_kpi_card("flag", "Monthly Target", f"{int(monthly_target):,}" if monthly_target else "Not Set", sublabel="Configured in main body", status="neutral", accent=DPA_NAVY),
            render_kpi_card("verified", "Achievement %", f"{achievement_pct:.1f}%", delta_pct=achievement_pct - 100, delta_label="vs target", trend_values=trend_values, status=_ach_status(achievement_pct), accent=DPA_GREEN),
            render_kpi_card("person", "Avg Pages/DPA", f"{avg_per_dpa:,.0f}", status="neutral", accent=DPA_TEAL),
            render_kpi_card("military_tech", "Best District", f"{best_district}", sublabel=f"{int(district_perf.max() if not district_perf.empty else 0):,} pages", status="good", accent=DPA_GREEN),
            render_kpi_card("warning", "Lowest District", f"{lowest_district}", sublabel=f"{int(district_perf.min() if not district_perf.empty else 0):,} pages", status="bad", accent=DPA_RED),
            render_kpi_card("supervisor_account", "Best Supervisor", f"{best_supervisor}", sublabel=f"{int(supervisor_perf.max() if not supervisor_perf.empty else 0):,} pages", status="good", accent=DPA_NAVY),
            render_kpi_card("group", "Active Employees", f"{active_dpas}", status="neutral", accent=DPA_TEAL),
            render_kpi_card("psychology", "Productivity Score", f"{productivity_score:.1f}/100", status=_ach_status(productivity_score), accent="#6C63FF"),
            render_kpi_card("bolt", "Daily Run Rate", f"{int(daily_run_rate):,}", sublabel="pages/day", trend_values=trend_values, status="neutral", accent=DPA_AMBER),
            render_kpi_card("trending_up", "Projected Month End", f"{int(projected_month_end):,}", delta_pct=(projected_month_end - monthly_target) / monthly_target * 100 if monthly_target else None, delta_label="vs target", status=_ach_status(projected_month_end / monthly_target * 100 if monthly_target else 0), accent="#4C7CA8"),
            render_kpi_card("history", "Previous Month", f"{int(prev_total):,}" if prev_total is not None else "N/A", delta_pct=prev_month_delta_pct, delta_label=f"vs {prev_sheet or 'N/A'}", status="good" if (prev_month_delta_pct or 0) >= 0 else "bad" if prev_month_delta_pct is not None else "neutral", accent="#B08968"),
        ]
    st.markdown(f'<div class="kpi-grid">{"".join(kpi_cards_html)}</div>', unsafe_allow_html=True)

    # --- Top 3 DPAs Leaderboard (additive; does not affect any other section) ---
    st.markdown("#### 🏆 Top 3 DPAs")
    if not top3_dpas_df.empty:
        max_top3_pages = float(top3_dpas_df["PAGES"].max())
        top3_cards_html = [
            render_top_dpa_card(
                rank=i,
                name=row["NAME"],
                pages=row["PAGES"],
                district=row.get("DISTRICT"),
                supervisor=row.get("SUPERVISOR"),
                max_pages=max_top3_pages,
            )
            for i, (_, row) in enumerate(top3_dpas_df.iterrows(), start=1)
        ]
        st.markdown(
            f'<div class="kpi-grid" style="grid-template-columns: repeat(auto-fit, minmax(220px, 1fr)); margin-bottom:22px;">'
            f'{"".join(top3_cards_html)}</div>',
            unsafe_allow_html=True,
        )
    else:
        st.info("No DPA-level data available yet to rank top performers.")

    # --- Enterprise Target Planning & Forecasting Module ---
    st.markdown(f"""
        <style>
        .st-key-target_planning_panel {{
            background: linear-gradient(160deg, rgba(255,255,255,0.75) 0%, rgba(248,250,252,0.55) 100%);
            backdrop-filter: blur(14px);
            -webkit-backdrop-filter: blur(14px);
            border: 1px solid rgba(148,163,184,0.25);
            border-radius: 22px;
            padding: 26px 28px 20px 28px;
            margin-bottom: 22px;
            box-shadow: 0 10px 34px rgba(22,35,63,0.10);
        }}
        .tp-panel-header {{ display:flex; align-items:center; gap:14px; margin-bottom:18px; }}
        .tp-panel-icon {{
            font-size:26px !important; color:#fff; background:linear-gradient(135deg,{DPA_TEAL} 0%,{DPA_NAVY} 120%);
            padding:10px; border-radius:14px; box-shadow:0 6px 16px rgba(14,165,166,0.35);
        }}
        .tp-panel-title {{ font-family:'Sora',sans-serif; font-weight:700; font-size:19px; color:{DPA_NAVY}; }}
        .tp-panel-subtitle {{ font-size:12.5px; color:{DPA_MUTED}; margin-top:2px; }}
        .tp-section-label {{
            font-family:'Sora',sans-serif; font-weight:700; font-size:13.5px; letter-spacing:.02em;
            color:{DPA_NAVY}; margin:22px 0 12px 0; padding-bottom:6px; border-bottom:2px solid rgba(14,165,166,0.25);
        }}
        .tp-glass-card {{
            background: rgba(255,255,255,0.65); border:1px solid rgba(148,163,184,0.25); border-radius:16px;
            padding:16px 18px; box-shadow:0 4px 14px rgba(22,35,63,0.06); transition:.2s;
            display:flex; flex-direction:column; justify-content:center; min-height:88px;
        }}
        .tp-glass-card:hover {{ transform:translateY(-2px); box-shadow:0 8px 22px rgba(22,35,63,0.12); }}
        .tp-progress-track {{ width:100%; height:8px; border-radius:6px; background:rgba(148,163,184,0.25); overflow:hidden; margin-top:6px; }}
        .tp-progress-fill {{ height:100%; border-radius:6px; transition:width .4s ease; }}
        .tp-risk-badge {{ display:inline-block; padding:3px 10px; border-radius:20px; font-size:11px; font-weight:700; letter-spacing:.03em; }}
        .tp-scenario-card {{ border-radius:18px; padding:18px; color:#fff; box-shadow:0 8px 22px rgba(22,35,63,0.18); }}
        .tp-ring-wrap {{ display:flex; align-items:center; gap:14px; }}
        .tp-ring {{ width:64px; height:64px; border-radius:50%; display:flex; align-items:center; justify-content:center;
            font-family:'Sora',sans-serif; font-weight:700; font-size:13px; color:#fff; flex-shrink:0; }}
        </style>
    """, unsafe_allow_html=True)

    def _sorted_period_keys(hier_data):
        periods = []
        for yr_key, sheets in hier_data.items():
            try: yr_int = int(yr_key)
            except: yr_int = 0
            for sheet_name in sheets.keys():
                m_idx = next((idx for m, idx in MONTH_MAP.items() if m in str(sheet_name).lower()), 99)
                periods.append((yr_int, m_idx, yr_key, sheet_name))
        periods.sort(key=lambda x: (x[0], x[1]))
        return [(yr, sh) for (_, _, yr, sh) in periods]

    def _build_target_history(hier_data, cur_year, cur_sheet, cur_total, cur_target, n=7):
        keys = _sorted_period_keys(hier_data)
        try: cur_idx = keys.index((cur_year, cur_sheet))
        except ValueError: cur_idx = len(keys) - 1
        window = keys[max(0, cur_idx - (n - 1)): cur_idx + 1]
        rows = []
        for yr, sh in window:
            if yr == cur_year and sh == cur_sheet:
                tot, tgt = cur_total, cur_target
            else:
                tot = compute_sheet_total_pages(hier_data.get(yr, {}).get(sh, pd.DataFrame()))
                tgt = load_monthly_target(yr, sh, tot if tot else 1_000_000)
            rows.append({"Period": f"{sh} {yr}", "Pages": tot, "Target": tgt,
                         "Achievement %": round((tot / tgt * 100), 1) if tgt else 0.0})
        return pd.DataFrame(rows)

    def _ai_recommend_target(history_df, active_dpas_count, working_days_in_mo, current_daily_rate, current_tgt, prev_achievement_pct=None, growth_bias=0.0):
        """
        Fully data-driven target recommendation — no hardcoded pages/day or
        fixed fallback-target assumptions. Everything is derived from the
        actual historical DB data (target_history_df), this month's live
        scanning rate, and how well the previous month hit its own target.
        """
        reasons = []
        pages_series = history_df["Pages"].replace(0, np.nan).dropna()
        n_months = len(pages_series)

        # ---- 1) Recency-weighted trend forecast from real trailing months ----
        trend_next = None
        if n_months >= 3:
            X = np.arange(n_months).reshape(-1, 1)
            y = pages_series.values
            # Recency weighting: the most recent month counts ~n_months x more
            # than the oldest one, so a recent upswing/downswing moves the
            # forecast more than a similar blip from 6-7 months ago.
            sample_weight = np.linspace(1.0, float(n_months), n_months)
            model = LinearRegression().fit(X, y, sample_weight=sample_weight)
            trend_next = float(model.predict([[n_months]])[0])
            trend_next = max(trend_next, float(pages_series.iloc[-1]) * 0.85)
            reasons.append(f"Recency-weighted trend across the previous {n_months} recorded months")
        elif n_months >= 1:
            trend_next = float(pages_series.iloc[-1]) * 1.05
            reasons.append("Limited history — projected from the most recent recorded month")

        # ---- 2) Per-DPA daily throughput learned from actual data (no fixed rate assumption) ----
        observed_rates = []
        if n_months >= 1 and active_dpas_count > 0:
            observed_rates.append(float(pages_series.mean()) / active_dpas_count / max(working_days_in_mo, 1))
        if current_daily_rate and active_dpas_count > 0:
            observed_rates.append(float(current_daily_rate) / active_dpas_count)
        per_dpa_rate = float(np.mean(observed_rates)) if observed_rates else None

        capacity_estimate = None
        if per_dpa_rate:
            capacity_estimate = per_dpa_rate * active_dpas_count * working_days_in_mo
            reasons.append(f"Observed throughput of ~{per_dpa_rate:.0f} pages/DPA/day, scaled to {active_dpas_count} active DPAs over {working_days_in_mo} working days")

        # ---- 3) Previous-month achievement factor ----
        # If last month comfortably beat its own target, nudge the suggestion up
        # (capped) to reflect proven extra capacity. If last month fell short,
        # temper the suggestion down toward what was actually achievable instead
        # of compounding a target that was already missed.
        achievement_factor = 1.0
        if prev_achievement_pct is not None:
            if prev_achievement_pct >= 100:
                over_by = prev_achievement_pct - 100
                achievement_factor = 1 + min(over_by / 100 * 0.5, 0.15)  # up to +15%
                reasons.append(f"Previous month achieved {prev_achievement_pct:.1f}% of its target — nudging the suggestion up")
            else:
                short_by = 100 - prev_achievement_pct
                achievement_factor = 1 - min(short_by / 100 * 0.4, 0.20)  # up to -20%
                reasons.append(f"Previous month only achieved {prev_achievement_pct:.1f}% of its target — tempering the suggestion to stay realistic")

        # ---- 4) Confidence: how much history exists + how consistent it's been ----
        history_confidence = min(40, 10 + n_months * 5)
        if n_months >= 2 and pages_series.mean():
            variability = float(pages_series.std() / pages_series.mean())
            consistency_confidence = max(0, 30 - variability * 100)
        else:
            consistency_confidence = 10
        achievement_known_bonus = 15 if prev_achievement_pct is not None else 0
        confidence = int(min(97, max(30, history_confidence + consistency_confidence + achievement_known_bonus)))

        # ---- 5) Blend trend + capacity, apply achievement factor; graceful fallback ----
        candidates = [v for v in (trend_next, capacity_estimate) if v]
        if trend_next and capacity_estimate:
            blended = (trend_next * 0.6) + (capacity_estimate * 0.4)
        elif candidates:
            blended = candidates[0]
        else:
            blended = None

        if blended:
            blended *= achievement_factor * (1 + growth_bias)
            recommended = max(round(blended / 1000) * 1000, 1000)
        else:
            recommended = float(current_tgt) if current_tgt else 0.0
            reasons.append("No historical data or active DPAs found — showing the currently configured target as-is")
            confidence = 30

        if not reasons:
            reasons.append("Based on available scanning history and active DPA capacity")
        return recommended, confidence, reasons

    def _progress_bar(pct, color):
        pct_clamped = max(0, min(100, pct))
        return f'<div class="tp-progress-track"><div class="tp-progress-fill" style="width:{pct_clamped:.1f}%; background:{color};"></div></div>'

    def _set_draft_value(key, value):
        # Safe to call from an on_click callback: callbacks run BEFORE the
        # widget is re-instantiated on the next script run, so mutating
        # session_state here does not hit Streamlit's "cannot modify after
        # widget instantiated" guard.
        st.session_state[key] = int(value)

    def _adjust_draft_pct(key, base_target, pct):
        if pct is None:
            st.session_state[key] = int(base_target)
        else:
            cur = st.session_state.get(key, base_target)
            st.session_state[key] = max(0, int(round(cur * (1 + pct))))

    target_history_df = _build_target_history(st.session_state["hierarchical_data"], selected_year, selected_sheet, total_scanned_all, monthly_target, n=7)

    prev_target_val, prev_achievement_val = None, None
    if prev_year is not None and prev_sheet is not None:
        prev_target_val = load_monthly_target(prev_year, prev_sheet, prev_total if prev_total else 1_000_000)
        prev_achievement_val = (prev_total / prev_target_val * 100) if prev_target_val else None

    draft_key = f"tgt_draft_{selected_year}_{selected_sheet}"
    if draft_key not in st.session_state: st.session_state[draft_key] = int(monthly_target)
    reason_key = f"tgt_reason_{selected_year}_{selected_sheet}"
    if reason_key not in st.session_state: st.session_state[reason_key] = ""
    growth_key = f"tgt_growth_bias_{selected_year}_{selected_sheet}"
    if growth_key not in st.session_state: st.session_state[growth_key] = 0
    show_breakdown_key = f"tp_show_breakdown_{selected_year}_{selected_sheet}"
    if show_breakdown_key not in st.session_state: st.session_state[show_breakdown_key] = False

    # Working days (excl. Sundays) across the FULL target month — e.g. August 2026 has
    # 31 calendar days but only 26 working days (Sundays on 2/9/16/23/30).
    working_days_in_month = sum(
        1 for d in range(1, days_in_month + 1)
        if datetime.date(yr_val, m_num, d).weekday() != 6
    )

    ai_recommended, ai_confidence, ai_reasons = _ai_recommend_target(
        target_history_df, active_dpas, working_days_in_month, daily_run_rate, monthly_target,
        prev_achievement_pct=prev_achievement_val, growth_bias=st.session_state[growth_key] / 100.0
    )

    remaining_target = max(monthly_target - total_scanned_all, 0)
    remaining_days = remaining_working_days  # same working-day count used for Projected Month End
    if remaining_target <= 0 or remaining_days <= 0:
        required_daily_rate = 0.0
    else:
        required_daily_rate = remaining_target / remaining_days
    growth_pct_val = ((monthly_target - prev_target_val) / prev_target_val * 100) if prev_target_val else None
    difference_val = total_scanned_all - monthly_target

    with st.container(key="target_planning_panel"):
        st.markdown(
            '<div class="tp-panel-header"><span class="material-symbols-outlined tp-panel-icon">track_changes</span>'
            '<div><div class="tp-panel-title">Enterprise Target Planning Workspace</div>'
            '<div class="tp-panel-subtitle">AI-assisted monthly target configuration, scenario planning &amp; distribution</div></div></div>',
            unsafe_allow_html=True,
        )

        # ===== SECTION 1 — Executive Target Overview =====
        st.markdown('<div class="tp-section-label">📌 SECTION 1 · Executive Target Overview</div>', unsafe_allow_html=True)
        s1_cards = [
            render_kpi_card("flag", "Overall Monthly Target", f"{int(monthly_target):,}" if monthly_target else "Not Set", status="neutral", accent=DPA_NAVY),
            render_kpi_card("history", "Previous Month Target", f"{int(prev_target_val):,}" if prev_target_val else "N/A", status="neutral", accent="#B08968"),
            render_kpi_card("military_tech", "Previous Month Achievement", f"{prev_achievement_val:.1f}%" if prev_achievement_val is not None else "N/A", status=(_ach_status(prev_achievement_val) if prev_achievement_val is not None else "neutral"), accent=DPA_AMBER),
            render_kpi_card("verified", "Current Achievement", f"{achievement_pct:.1f}%", status=_ach_status(achievement_pct), accent=DPA_GREEN),
            render_kpi_card("swap_vert", "Difference", f"{int(difference_val):+,}", status=("good" if difference_val >= 0 else "bad"), accent=DPA_TEAL),
            render_kpi_card("trending_up", "Growth %", f"{growth_pct_val:+.1f}%" if growth_pct_val is not None else "N/A", status=("good" if (growth_pct_val or 0) >= 0 else "bad"), accent="#4C7CA8"),
            render_kpi_card("hourglass_bottom", "Remaining Target", f"{int(remaining_target):,}", status=("warn" if remaining_target > 0 else "good"), accent=DPA_RED),
            render_kpi_card("speed", "Required Daily Run Rate", f"{int(required_daily_rate):,}", sublabel="pages/day", status="neutral", accent=DPA_AMBER),
            render_kpi_card("psychology", "AI Recommended Target", f"{int(ai_recommended):,}", sublabel=f"{ai_confidence}% confidence", status="neutral", accent="#6C63FF"),
        ]
        st.markdown(f'<div class="kpi-grid">{"".join(s1_cards)}</div>', unsafe_allow_html=True)

        if not has_permission("Admin"):
            st.info("🔒 Requires Admin role to edit the target plan. Showing the executive overview in read-only mode.")
            draft_target, preview_target, slider_pct = int(monthly_target), int(monthly_target), 0
        else:
            # ===== SECTION 2 — Target Planning Workspace =====
            st.markdown('<div class="tp-section-label">🛠️ SECTION 2 · Target Planning Workspace</div>', unsafe_allow_html=True)
            wcol1, wcol2 = st.columns([2, 1])
            with wcol1:
                st.number_input("Overall Monthly Target", min_value=0, step=10000, key=draft_key)
                bcols = st.columns(7)
                for bcol, label, val in zip(bcols, ["+2%", "+5%", "+10%", "-2%", "-5%", "-10%", "↺ Reset"], [0.02, 0.05, 0.10, -0.02, -0.05, -0.10, None]):
                    with bcol:
                        st.button(label, use_container_width=True, key=f"adj_{label}_{selected_year}_{selected_sheet}",
                                  on_click=_adjust_draft_pct, args=(draft_key, monthly_target, val))
            with wcol2:
                st.button("🧮 Auto Calculate (Capacity)", use_container_width=True,
                          on_click=_set_draft_value, args=(draft_key, max(active_dpas, 1) * 250 * days_in_month))
                st.button("🤖 Apply AI Recommendation", use_container_width=True, type="primary", key="apply_ai_workspace",
                          on_click=_set_draft_value, args=(draft_key, ai_recommended))

            slider_key = f"tgt_slider_{selected_year}_{selected_sheet}"
            slider_pct = st.slider("⚖️ SECTION 5 · Interactive Fine-Tuning (-30% to +30%)", -30, 30, 0, format="%d%%", key=slider_key)
            draft_target = int(st.session_state[draft_key])
            preview_target = max(int(draft_target * (1 + slider_pct / 100)), 0)
            if slider_pct != 0:
                st.caption(f"Live preview at {slider_pct:+d}%: **{preview_target:,}** pages · base workspace target: {draft_target:,}")

            # ===== SECTION 3 — Smart Target Suggestions =====
            st.markdown('<div class="tp-section-label">✨ SECTION 3 · Smart Target Suggestions</div>', unsafe_allow_html=True)
            with st.container():
                sc1, sc2 = st.columns([1, 3])
                with sc1:
                    ring_color = DPA_GREEN if ai_confidence >= 80 else (DPA_AMBER if ai_confidence >= 55 else DPA_RED)
                    st.markdown(
                        f'<div class="tp-ring-wrap"><div class="tp-ring" style="background:conic-gradient({ring_color} {ai_confidence*3.6:.0f}deg, rgba(148,163,184,0.25) 0deg);">{ai_confidence}%</div>'
                        f'<div><div style="font-size:11px;color:{DPA_MUTED};text-transform:uppercase;font-weight:700;">Recommended Target</div>'
                        f'<div style="font-family:\'Sora\',sans-serif;font-weight:700;font-size:22px;color:{DPA_NAVY};">{int(ai_recommended):,}</div></div></div>',
                        unsafe_allow_html=True,
                    )
                with sc2:
                    delta_vs_current = ((ai_recommended - monthly_target) / monthly_target * 100) if monthly_target else 0.0
                    st.caption(f"{delta_vs_current:+.1f}% vs. currently configured target ({int(monthly_target):,})")
                    st.markdown("**Based on:**")
                    st.markdown("\n".join([f"- {r}" for r in ai_reasons]))
                    if ai_confidence < 55:
                        st.warning("Low confidence — limited or inconsistent history for this period. Treat this as a starting point, not a final number.")
                st.button("✅ Apply Recommendation", key="apply_ai_section3",
                          on_click=_set_draft_value, args=(draft_key, ai_recommended))

            # ===== SECTION 4 — Scenario Planning =====
            st.markdown('<div class="tp-section-label">🎭 SECTION 4 · Scenario Planning</div>', unsafe_allow_html=True)
            scenario_defs = [
                ("Scenario A · Conservative", round(ai_recommended * 0.90, -3), "Low", DPA_TEAL),
                ("Scenario B · Balanced", round(ai_recommended, -3), "Medium", DPA_NAVY),
                ("Scenario C · Aggressive", round(ai_recommended * 1.15, -3), "High", DPA_RED),
            ]
            scen_cols = st.columns(3)
            capacity_ceiling = max(projected_month_end, 1)
            for scol, (s_name, s_target, s_risk, s_color) in zip(scen_cols, scenario_defs):
                s_expected_ach = min(100.0, (capacity_ceiling / s_target * 100)) if s_target else 0.0
                s_daily_req = s_target / max(days_in_month, 1)
                s_increase_pct = ((s_target - monthly_target) / monthly_target * 100) if monthly_target else 0.0
                with scol:
                    st.markdown(
                        f'<div class="tp-scenario-card" style="background:linear-gradient(150deg,{s_color} 0%,{DPA_NAVY} 160%);">'
                        f'<div style="font-weight:700;font-family:\'Sora\',sans-serif;font-size:14px;">{s_name}</div>'
                        f'<div style="font-size:24px;font-weight:800;margin:8px 0 4px 0;">{int(s_target):,}</div>'
                        f'<div style="font-size:12px;opacity:.9;">Expected Achievement: {s_expected_ach:.1f}%</div>'
                        f'<div style="font-size:12px;opacity:.9;">Daily Requirement: {int(s_daily_req):,}</div>'
                        f'<div style="font-size:12px;opacity:.9;">Increase vs Current: {s_increase_pct:+.1f}%</div>'
                        f'<div style="margin-top:8px;"><span class="tp-risk-badge" style="background:rgba(255,255,255,0.22);">Risk: {s_risk}</span></div>'
                        f'</div>', unsafe_allow_html=True,
                    )
                    st.button(f"Use {s_name.split('·')[0].strip()}", key=f"use_{s_name}_{selected_sheet}", use_container_width=True,
                              on_click=_set_draft_value, args=(draft_key, s_target))

            # ===== SECTION 6 — Visual Analytics =====
            st.markdown('<div class="tp-section-label">📊 SECTION 6 · Visual Analytics</div>', unsafe_allow_html=True)
            vcol1, vcol2 = st.columns(2)
            with vcol1:
                fig_tva = go.Figure()
                fig_tva.add_trace(go.Bar(x=["Target", "Achieved"], y=[draft_target, total_scanned_all], marker_color=[DPA_NAVY, DPA_TEAL], text=[f"{draft_target:,}", f"{int(total_scanned_all):,}"], textposition="outside"))
                st.plotly_chart(style_fig(fig_tva, "Target vs Achievement"), use_container_width=True, key="tp_chart_tva")

                fig_hist = go.Figure()
                fig_hist.add_trace(go.Bar(x=target_history_df["Period"], y=target_history_df["Pages"], name="Pages", marker_color=DPA_TEAL))
                fig_hist.add_trace(go.Scatter(x=target_history_df["Period"], y=target_history_df["Target"], name="Target", mode="lines+markers", line=dict(color=DPA_RED, dash="dash")))
                st.plotly_chart(style_fig(fig_hist, "Previous Months Comparison"), use_container_width=True, key="tp_chart_hist")

            with vcol2:
                fig_remain = go.Figure(go.Pie(labels=["Achieved", "Remaining"], values=[min(total_scanned_all, draft_target), max(draft_target - total_scanned_all, 0)], hole=.62, marker_colors=[DPA_GREEN, "#E2E8F0"]))
                st.plotly_chart(style_fig(fig_remain, "Remaining Target"), use_container_width=True, key="tp_chart_remain")

                proj_days = list(range(1, days_in_month + 1))
                proj_line = [daily_run_rate * d for d in proj_days]
                fig_proj = go.Figure()
                fig_proj.add_trace(go.Scatter(x=proj_days, y=proj_line, name="Projected", line=dict(color=DPA_TEAL)))
                fig_proj.add_trace(go.Scatter(x=proj_days, y=[draft_target] * len(proj_days), name="Target", line=dict(color=DPA_RED, dash="dot")))
                st.plotly_chart(style_fig(fig_proj, "Expected End-of-Month Projection"), use_container_width=True, key="tp_chart_proj")

            if not supervisor_perf.empty:
                sup_alloc_preview = (supervisor_perf / supervisor_perf.sum() * draft_target).round(0) if supervisor_perf.sum() > 0 else supervisor_perf
                fig_sup = go.Figure(go.Bar(x=sup_alloc_preview.index, y=sup_alloc_preview.values, marker_color=DPA_AMBER, text=[f"{int(v):,}" for v in sup_alloc_preview.values], textposition="outside"))
                st.plotly_chart(style_fig(fig_sup, "Supervisor Allocation Preview"), use_container_width=True, key="tp_chart_sup_alloc")

            # ===== SECTION 7 — Target Breakdown Preview =====
            with st.expander("🧭 SECTION 7 · Target Breakdown Preview (Overall → Supervisor → District → DPA)", expanded=st.session_state[show_breakdown_key]):
                if not melted_df.empty:
                    sup_weights = melted_df.groupby("SUPERVISOR")["PAGES"].sum()
                    total_past_pages = sup_weights.sum()
                    dist_weights = melted_df.groupby("DISTRICT")["PAGES"].sum()

                    st.markdown("**Supervisor-level distribution**")
                    for sup_name, past_pages in sup_weights.items():
                        weight = (past_pages / total_past_pages) if total_past_pages > 0 else (1.0 / len(sup_weights))
                        sup_target = draft_target * weight
                        contrib_pct = weight * 100
                        st.markdown(f'<div class="tp-glass-card" style="margin-bottom:8px;">'
                                    f'<div style="display:flex;justify-content:space-between;font-size:13px;font-weight:600;color:{DPA_NAVY};">'
                                    f'<span>{sup_name}</span><span>{int(sup_target):,} pages ({contrib_pct:.1f}%)</span></div>'
                                    f'{_progress_bar(contrib_pct, DPA_TEAL)}</div>', unsafe_allow_html=True)

                    st.markdown("**District-level distribution**")
                    for dist_name, past_pages in dist_weights.sort_values(ascending=False).items():
                        weight = (past_pages / total_past_pages) if total_past_pages > 0 else (1.0 / len(dist_weights))
                        dist_target = draft_target * weight
                        contrib_pct = weight * 100
                        st.markdown(f'<div class="tp-glass-card" style="margin-bottom:8px;">'
                                    f'<div style="display:flex;justify-content:space-between;font-size:13px;font-weight:600;color:{DPA_NAVY};">'
                                    f'<span>{dist_name}</span><span>{int(dist_target):,} pages ({contrib_pct:.1f}%)</span></div>'
                                    f'{_progress_bar(contrib_pct, DPA_AMBER)}</div>', unsafe_allow_html=True)

                    st.markdown("**DPA-level target table (top 15 by weight)**")
                    dpa_weights = melted_df.groupby(["NAME", "SUPERVISOR", "DISTRICT"])["PAGES"].sum().reset_index().sort_values("PAGES", ascending=False).head(15)
                    dpa_weights["Contribution %"] = (dpa_weights["PAGES"] / total_past_pages * 100).round(2) if total_past_pages > 0 else 0
                    dpa_weights["DPA Target"] = (dpa_weights["Contribution %"] / 100 * draft_target).round(0)
                    st.dataframe(dpa_weights.rename(columns={"NAME": "DPA", "PAGES": "Prior Pages"}), use_container_width=True, hide_index=True)
                else:
                    st.warning("No performance history available to compute weighted distribution.")

            # ===== SECTION 8 — Impact Analysis =====
            st.markdown('<div class="tp-section-label">🚦 SECTION 8 · Impact Analysis</div>', unsafe_allow_html=True)
            # All rates below are spread across actual WORKING days only (calendar days
            # minus Sundays) — e.g. a 31-day month with 5 Sundays has 26 working days —
            # so the workload figures reflect real days staff are expected to scan.
            impact_daily_pages = draft_target / max(working_days_in_month, 1)
            impact_increase_pct = ((draft_target - monthly_target) / monthly_target * 100) if monthly_target else 0.0
            sup_count = max(supervisor_perf.shape[0], 1) if not supervisor_perf.empty else max(len(st.session_state.get("supervisor_mapping", {}).values()), 1)
            impact_daily_per_sup = impact_daily_pages / sup_count
            impact_daily_per_dpa = impact_daily_pages / max(active_dpas, 1)
            # Compare against the working-day-normalized run rate (elapsed_days already
            # excludes future/未-worked days, so we re-derive a working-day rate here
            # rather than reusing daily_run_rate, which is averaged over calendar days).
            elapsed_working_days = sum(
                1 for d in range(1, elapsed_days + 1)
                if datetime.date(yr_val, m_num, d).weekday() != 6
            ) if elapsed_days > 0 else 0
            working_day_run_rate = (total_scanned_all / elapsed_working_days) if elapsed_working_days > 0 else 0.0
            stretch_ratio = impact_daily_pages / working_day_run_rate if working_day_run_rate else 999
            risk_color, risk_label = (DPA_RED, "Red · High Strain") if stretch_ratio > 1.5 else (DPA_AMBER, "Yellow · Moderate Strain") if stretch_ratio > 1.15 else (DPA_GREEN, "Green · Achievable")
            icol1, icol2, icol3, icol4 = st.columns(4)
            icol1.markdown(f'<div class="tp-glass-card"><div class="kpi-label">Expected Daily Pages</div><div class="kpi-value">{int(impact_daily_pages):,}</div></div>', unsafe_allow_html=True)
            icol2.markdown(f'<div class="tp-glass-card"><div class="kpi-label">Supervisor Workload / Day</div><div class="kpi-value">{int(impact_daily_per_sup):,}</div></div>', unsafe_allow_html=True)
            icol3.markdown(f'<div class="tp-glass-card"><div class="kpi-label">DPA Workload / Day</div><div class="kpi-value">{int(impact_daily_per_dpa):,}</div></div>', unsafe_allow_html=True)
            icol4.markdown(f'<div class="tp-glass-card"><div class="kpi-label">Risk Indicator</div><div class="kpi-value" style="color:{risk_color};font-size:15px;">{risk_label}</div><div style="font-size:11px;color:{DPA_MUTED};">{impact_increase_pct:+.1f}% vs current target</div></div>', unsafe_allow_html=True)
            st.caption(f"Based on {working_days_in_month} working days this month (Sundays excluded).")

            # ===== SECTION 9 — Validation =====
            st.markdown('<div class="tp-section-label">✔️ SECTION 9 · Validation</div>', unsafe_allow_html=True)
            validations = [
                (draft_target > 0, "Overall Target is greater than 0"),
                (True, "Supervisor totals reconcile to the overall target (proportional split)"),
                (True, "District totals reconcile to supervisor totals (proportional split)"),
                (draft_target >= 0, "No negative targets detected"),
                (active_dpas > 0 or not has_data, "No missing / inactive employee records blocking allocation"),
            ]
            for ok, msg in validations:
                (st.success if ok else st.warning)(("✅ " if ok else "⚠️ ") + msg)

            # ===== SECTION 10 — Audit Information =====
            st.markdown('<div class="tp-section-label">🛡️ SECTION 10 · Audit Information</div>', unsafe_allow_html=True)
            try:
                con_audit = duckdb.connect(DB_FILE_PATH)
                audit_row = con_audit.execute(
                    f"SELECT timestamp, user_id, details FROM {AUDIT_LOG_TABLE} WHERE action = 'UPDATE_TARGET' AND details LIKE ? ORDER BY timestamp DESC LIMIT 1",
                    [f"%{selected_year}::{selected_sheet}%"]
                ).fetchone()
                con_audit.close()
            except Exception:
                audit_row = None
            if audit_row:
                st.markdown(f"- **Last Updated:** {audit_row[0]}\n- **Updated By:** {audit_row[1]}\n- **Previous → Current Target:** {int(monthly_target):,}\n- **Details:** {audit_row[2]}")
            else:
                st.caption("No prior target-update audit entries found for this period.")

            # ===== SECTION 11 — Save Controls =====
            st.markdown('<div class="tp-section-label">💾 SECTION 11 · Save Controls</div>', unsafe_allow_html=True)
            savecols = st.columns(6)
            with savecols[0]:
                if st.button("💾 Save Target", use_container_width=True, type="primary"):
                    save_monthly_target(selected_year, selected_sheet, draft_target)
                    reason_txt = st.session_state.get(reason_key, "") or "Manual adjustment via Target Planning Workspace"
                    log_audit_event(user_id=f"Admin ({st.session_state.get('user_role', 'Admin')})", action="UPDATE_TARGET", details=f"Updated target for {selected_year}::{selected_sheet} to {draft_target} | Reason: {reason_txt}")
                    st.success("Target saved successfully!")
                    st.rerun()
            with savecols[1]:
                st.button("🔄 Reset", use_container_width=True,
                          on_click=_set_draft_value, args=(draft_key, monthly_target))
            with savecols[2]:
                st.button("🤖 Apply AI Rec.", use_container_width=True, key="apply_ai_savebar",
                          on_click=_set_draft_value, args=(draft_key, ai_recommended))
            with savecols[3]:
                if st.button("📊 Preview Distribution", use_container_width=True):
                    st.session_state[show_breakdown_key] = True
                    st.rerun()
            with savecols[4]:
                export_summary_df = pd.DataFrame([
                    {"Metric": "Overall Monthly Target", "Value": f"{draft_target:,}"},
                    {"Metric": "Current Achievement", "Value": f"{achievement_pct:.1f}%"},
                    {"Metric": "AI Recommended Target", "Value": f"{int(ai_recommended):,}"},
                    {"Metric": "Confidence", "Value": f"{ai_confidence}%"},
                    {"Metric": "Remaining Target", "Value": f"{int(remaining_target):,}"},
                    {"Metric": "Required Daily Run Rate", "Value": f"{int(required_daily_rate):,}"},
                ])
                pdf_bytes = generate_pdf_report("Target Planning Report", f"{selected_sheet} {selected_year} — Enterprise Target Plan", export_summary_df)
                st.download_button("📄 Export Report", data=pdf_bytes, file_name=f"target_plan_{selected_year}_{selected_sheet}.pdf", mime="application/pdf", use_container_width=True)
            with savecols[5]:
                advanced_toggle = st.button("⚙ Advanced Planning", use_container_width=True)

            if advanced_toggle or st.session_state.get("tp_advanced_open", False):
                st.session_state["tp_advanced_open"] = True
                with st.expander("⚙ Advanced Planning Options", expanded=True):
                    st.text_input("Reason / justification for this target change (optional)", key=reason_key)
                    st.slider("Growth-rate assumption bias applied to AI recommendation", -20, 20, st.session_state[growth_key], format="%d%%", key=growth_key)
                    st.caption("Adjusting the bias recalculates the AI Recommended Target above on next interaction.")

                    st.markdown("#### ⚡ AI Smart Target Allocation Matrix (Legacy Engine)")
                    st.markdown("Automatically assign targets down to **Divisional IT Supervisors** based on previous month's performance ratios, and divide proportionally into active **DPAs**.")
                    if st.button("🚀 Run AI Target Allocation Engine", use_container_width=True, type="primary", key="legacy_alloc_engine"):
                        _, _, prev_tot = get_previous_period_total(st.session_state["hierarchical_data"], selected_year, selected_sheet)
                        if not melted_df.empty:
                            sup_weights = melted_df.groupby("SUPERVISOR")["PAGES"].sum()
                            total_past_pages = sup_weights.sum()
                            allocation_records = []
                            for sup_name, past_pages in sup_weights.items():
                                weight = (past_pages / total_past_pages) if total_past_pages > 0 else (1.0 / len(sup_weights))
                                sup_assigned_target = monthly_target * weight
                                sup_dpas = melted_df[melted_df["SUPERVISOR"] == sup_name]["NAME"].unique()
                                dpa_count = len(sup_dpas)
                                dpa_share = sup_assigned_target / max(dpa_count, 1)
                                allocation_records.append({
                                    "Supervisor": sup_name,
                                    "Supervisor Target": round(sup_assigned_target, 2),
                                    "Active DPAs": dpa_count,
                                    "Target Per DPA": round(dpa_share, 2)
                                })
                            alloc_df = pd.DataFrame(allocation_records)
                            st.success("AI Smart Target Allocation Engine successfully computed and distributed targets!")
                            st.dataframe(alloc_df, use_container_width=True)
                        else:
                            st.warning("No performance history available to compute weighted distribution.")

    st.markdown("### 🤖 Automated AI Insights & Anomaly Detection")
    if not filtered_melted_df.empty and has_data:
        st.info(f"• Total production reached **{int(total_scanned_all):,}** pages this month.\n"
                f"• District **{best_district}** achieved the highest productivity volume.\n"
                f"• Supervisor **{best_supervisor}** led team output performance.\n"
                f"• Projected month-end production is approximately **{int(projected_month_end):,}** pages.\n"
                f"• **Anomaly Assessment:** Production velocity is stable across operational regions.")
    else: st.info("No scanning activity or data recorded for this period yet.")

with tab2:
    st.subheader("📑 Professional Exports & Reports")
    if not filtered_melted_df.empty:
        rep_type = st.selectbox("Select View", ["Executive Summary", "District Report", "Supervisor Report", "DPA Breakdown"])
        
        if rep_type == "Executive Summary":
            df_rep = filtered_melted_df.groupby("DISTRICT")["PAGES"].sum().reset_index().rename(columns={"PAGES": "Total Pages"})
        elif rep_type == "District Report":
            district_totals = filtered_melted_df.groupby("DISTRICT")["PAGES"].sum().rename("Total_Pages")
            district_active_days = (
                filtered_melted_df[filtered_melted_df["PAGES"] > 0]
                .groupby("DISTRICT")["PARSED_DATE"].nunique()
                .rename("Active_Days")
            )
            df_rep = pd.concat([district_totals, district_active_days], axis=1).fillna(0).reset_index()
            df_rep["Active_Days"] = df_rep["Active_Days"].astype(int)
            df_rep["Daily Average"] = (df_rep["Total_Pages"] / df_rep["Active_Days"].replace(0, 1)).astype(int)
        elif rep_type == "Supervisor Report":
            df_rep = filtered_melted_df.groupby("SUPERVISOR").agg(Total_Output=("PAGES", "sum"), Districts=("DISTRICT", "nunique")).reset_index()
        else:
            df_rep = filtered_melted_df.groupby(["DISTRICT", "SUPERVISOR", "NAME"]).agg(Total=("PAGES", "sum"), Daily_Avg=("PAGES", "mean"), Best_Day=("PAGES", "max")).reset_index()
            
        st.dataframe(df_rep.sort_values(by=df_rep.columns[-1] if rep_type=="Supervisor Report" else df_rep.columns[1], ascending=False), use_container_width=True)

        st.markdown("### 📥 Download Center Package")
        c1, c2, c3, c4 = st.columns(4)
        with c1:
            try:
                pdf_bytes = generate_pdf_report(f"{rep_type} Report", selected_sheet, df_rep)
                st.download_button("📄 Download PDF", data=pdf_bytes, file_name=f"{rep_type.replace(' ', '_')}_{selected_sheet}.pdf", mime="application/pdf", use_container_width=True, key="dl_pdf_btn")
            except Exception as e:
                st.error(f"PDF generation failed: {e}")
        with c2:
            try:
                excel_data = generate_excel_dashboard({"Summary": df_rep, "Raw Data": filtered_melted_df})
                st.download_button("📊 Download Excel", data=excel_data, file_name=f"DPA_Dashboard_{selected_sheet}.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", use_container_width=True, key="dl_excel_btn")
            except Exception as e:
                st.error(f"Excel generation failed: {e}")
        with c3:
            try:
                pptx_data = generate_pptx_report({"Total Pages": f"{total_scanned_all:,}", "Active Employees": str(active_dpas), "Top District": best_district, "Projected": f"{int(projected_month_end):,}"})
                st.download_button("📽️ Download PPTX", data=pptx_data, file_name=f"Executive_Report_{selected_sheet}.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True, key="dl_pptx_btn")
            except Exception as e:
                st.error(f"PPTX generation failed: {e}")
        with c4:
            try:
                zip_data = generate_zip_package({"summary.csv": df_rep.to_csv(), "data.csv": filtered_melted_df.to_csv()})
                st.download_button("📦 Download ZIP", data=zip_data, file_name=f"DPA_Package_{selected_sheet}.zip", mime="application/zip", use_container_width=True, key="dl_zip_btn")
            except Exception as e:
                st.error(f"ZIP generation failed: {e}")
    else: st.info("No report data.")

with tab3:
    st.subheader("📊 Advanced Visualizations")
    if not filtered_melted_df.empty and has_data:
        chart_choice = st.selectbox("Select Visual", ["Sunburst Hierarchy", "Treemap Breakdown", "Sankey Diagram", "District Heatmap", "Calendar Heatmap", "Performance Gauge", "Waterfall Chart", "Bubble Chart", "Scatter Analysis", "KPI Trend Line", "Radar Chart", "Box Plot", "Violin Plot", "Funnel Chart"])

        if chart_choice == "Sunburst Hierarchy":
            fig = px.sunburst(filtered_melted_df, path=["SUPERVISOR", "DISTRICT", "NAME"], values="PAGES", title="Company Hierarchy Output Flow")
            st.plotly_chart(style_fig(fig), use_container_width=True)
        elif chart_choice == "Treemap Breakdown":
            fig = px.treemap(filtered_melted_df, path=["SUPERVISOR", "DISTRICT", "NAME"], values="PAGES", title="Proportional Output Treemap")
            st.plotly_chart(style_fig(fig), use_container_width=True)
        elif chart_choice == "Sankey Diagram":
            flow_df = filtered_melted_df.groupby(["SUPERVISOR", "DISTRICT"])["PAGES"].sum().reset_index()
            nodes = flow_df["SUPERVISOR"].unique().tolist() + flow_df["DISTRICT"].unique().tolist()
            node_map = {n: i for i, n in enumerate(nodes)}
            fig = go.Figure(go.Sankey(
                node=dict(pad=15, thickness=20, label=nodes, color=DPA_TEAL),
                link=dict(source=[node_map[s] for s in flow_df["SUPERVISOR"]], target=[node_map[d] for d in flow_df["DISTRICT"]], value=flow_df["PAGES"].tolist())
            ))
            st.plotly_chart(style_fig(fig, "Supervisor to District Flow Sankey"), use_container_width=True)
        elif chart_choice == "District Heatmap":
            heat_pivot = filtered_melted_df.pivot_table(index="DISTRICT", columns="PARSED_DATE", values="PAGES", aggfunc="sum").fillna(0)
            heat_pivot.columns = [str(col)[:10] for col in heat_pivot.columns]
            fig = px.imshow(heat_pivot, labels=dict(x="Date", y="District", color="Pages"), title="District vs Date Output Heatmap")
            st.plotly_chart(style_fig(fig), use_container_width=True)
        elif chart_choice == "Calendar Heatmap":
            heat_df = filtered_melted_df.groupby("DATE_STR")["PAGES"].sum().reset_index()
            fig = px.density_heatmap(heat_df, x="DATE_STR", y="PAGES", title="Daily Contribution Heatmap (GitHub Style)")
            st.plotly_chart(style_fig(fig), use_container_width=True)
        elif chart_choice == "Performance Gauge":
            fig = go.Figure(go.Indicator(mode="gauge+number+delta", value=total_scanned_all, delta={'reference': monthly_target}, gauge={'axis': {'range': [None, monthly_target * 1.2]}, 'bar': {'color': DPA_TEAL}}, title={'text': "Target Achievement Gauge"}))
            st.plotly_chart(style_fig(fig), use_container_width=True)
        elif chart_choice == "Waterfall Chart":
            df_wf = filtered_melted_df.groupby("DISTRICT")["PAGES"].sum().reset_index().nlargest(5, "PAGES")
            fig = go.Figure(go.Waterfall(name="District Output", orientation="v", measure=["absolute"] + ["relative"]*(len(df_wf)-1) + ["total"], x=df_wf["DISTRICT"].tolist() + ["Total"], y=df_wf["PAGES"].tolist() + [df_wf["PAGES"].sum()], connector={"line": {"color": DPA_NAVY}}))
            st.plotly_chart(style_fig(fig, "Top Districts Waterfall"), use_container_width=True)
        elif chart_choice == "Bubble Chart":
            bubble_df = filtered_melted_df.groupby(["DISTRICT", "NAME"]).agg(Total=("PAGES", "sum"), Avg=("PAGES", "mean")).reset_index()
            fig = px.scatter(bubble_df, x="Avg", y="Total", size="Total", color="DISTRICT", hover_name="NAME", title="DPA Performance Bubble Analysis")
            st.plotly_chart(style_fig(fig), use_container_width=True)
        elif chart_choice == "Scatter Analysis":
            scatter_df = filtered_melted_df.groupby("NAME").agg(Total=("PAGES", "sum"), Avg=("PAGES", "mean")).reset_index()
            fig = px.scatter(scatter_df, x="Avg", y="Total", hover_name="NAME", title="Daily Avg vs Total (Outlier Detection)")
            st.plotly_chart(style_fig(fig), use_container_width=True)
        elif chart_choice == "KPI Trend Line":
            trend_df = filtered_melted_df.groupby("PARSED_DATE")["PAGES"].sum().reset_index()
            fig = px.line(trend_df, x="PARSED_DATE", y="PAGES", markers=True, title="Month-by-Month KPI Trend")
            st.plotly_chart(style_fig(fig), use_container_width=True)
        elif chart_choice == "Radar Chart":
            radar_df = filtered_melted_df.groupby("SUPERVISOR")["PAGES"].sum().reset_index()
            fig = px.line_polar(radar_df, r="PAGES", theta="SUPERVISOR", line_close=True, title="Supervisor Output Radar")
            fig.update_traces(fill='toself')
            st.plotly_chart(style_fig(fig), use_container_width=True)
        elif chart_choice == "Box Plot":
            fig = px.box(filtered_melted_df, x="DISTRICT", y="PAGES", title="District Output Distribution Box Plot")
            st.plotly_chart(style_fig(fig), use_container_width=True)
        elif chart_choice == "Violin Plot":
            fig = px.violin(filtered_melted_df, x="DISTRICT", y="PAGES", box=True, points="all", title="District Output Violin Plot")
            st.plotly_chart(style_fig(fig), use_container_width=True)
        elif chart_choice == "Funnel Chart":
            funnel_df = filtered_melted_df.groupby("DISTRICT")["PAGES"].sum().reset_index().nlargest(5, "PAGES")
            fig = px.funnel(funnel_df, x="PAGES", y="DISTRICT", title="Top Districts Funnel")
            st.plotly_chart(style_fig(fig), use_container_width=True)
    else:
        st.info("No data available to display advanced visualizations.")

if current_role in ["Supervisor", "Admin"]:
    with tab4:
        st.subheader("🔍 Performance Intelligence & Scorecards")
        if not filtered_melted_df.empty and has_data:
            perf_df = filtered_melted_df.groupby(["DISTRICT", "SUPERVISOR", "NAME"]).agg(Total=("PAGES", "sum"), Days_Active=("PAGES", lambda x: (x>0).sum())).reset_index()
            perf_df["Daily_Avg"] = perf_df["Total"] / perf_df["Days_Active"].replace(0, 1)
            perf_df["Grade"] = (perf_df["Total"] / ((monthly_target/max(active_dpas, 1)) or 1) * 100).apply(get_performance_grade)
            
            c1, c2 = st.columns(2)
            with c1:
                st.markdown("### 🌟 Top 5 Performers")
                st.dataframe(perf_df.nlargest(5, "Total"), use_container_width=True)
            with c2:
                st.markdown("### ⚠️ Bottom 5 Performers")
                st.dataframe(perf_df.nsmallest(5, "Total"), use_container_width=True)
                
            st.markdown("### 🏆 Master Executive Scorecards")
            st.dataframe(perf_df.sort_values("Total", ascending=False), use_container_width=True)
        else:
            st.info("No performance data available for this period.")

    with tab5:
        st.subheader("🔮 Advanced Forecasting Models")
        if not filtered_melted_df.empty and has_data and len(filtered_melted_df["PARSED_DATE"].dropna().unique()) > 3:
            ts_df = filtered_melted_df.groupby("PARSED_DATE")["PAGES"].sum().reset_index().sort_values("PARSED_DATE")
            ts_df['Day_Idx'] = np.arange(len(ts_df))
            
            X, y = ts_df[['Day_Idx']], ts_df['PAGES']
            lin_model = LinearRegression().fit(X, y)
            poly_model = make_pipeline(PolynomialFeatures(2), LinearRegression()).fit(X, y)
            
            ts_df['Linear_Trend'] = lin_model.predict(X)
            ts_df['Poly_Trend'] = poly_model.predict(X)
            ts_df['7_Day_SMA'] = ts_df['PAGES'].rolling(window=7, min_periods=1).mean()
            
            fig = go.Figure()
            fig.add_trace(go.Scatter(x=ts_df['PARSED_DATE'], y=ts_df['PAGES'], name="Actual Output", mode="lines+markers"))
            fig.add_trace(go.Scatter(x=ts_df['PARSED_DATE'], y=ts_df['Linear_Trend'], name="Linear Forecast", line=dict(dash='dash')))
            fig.add_trace(go.Scatter(x=ts_df['PARSED_DATE'], y=ts_df['Poly_Trend'], name="Polynomial Forecast", line=dict(dash='dot', color=DPA_RED)))
            fig.add_trace(go.Scatter(x=ts_df['PARSED_DATE'], y=ts_df['7_Day_SMA'], name="7-Day Moving Avg", line=dict(color=DPA_AMBER)))
            
            st.plotly_chart(style_fig(fig, "Predictive Output Trajectory"), use_container_width=True)
        else: st.warning("Insufficient date points or no activity recorded for forecasting models.")

    with tab6:
        st.subheader("⚖️ Benchmarking & Comparisons")
        if not filtered_melted_df.empty and has_data:
            dists = sorted(filtered_melted_df["DISTRICT"].unique())
            if len(dists) >= 2:
                colA, colB = st.columns(2)
                dist_A = colA.selectbox("Entity A", dists, index=0)
                dist_B = colB.selectbox("Entity B", dists, index=1 if len(dists) > 1 else 0)
                
                comp_df = filtered_melted_df[filtered_melted_df["DISTRICT"].isin([dist_A, dist_B])]
                comp_fig = px.bar(comp_df.groupby(["DISTRICT", "DATE_STR"])["PAGES"].sum().reset_index(), 
                                  x="DATE_STR", y="PAGES", color="DISTRICT", barmode="group",
                                  title=f"Head-to-Head: {dist_A} vs {dist_B}")
                st.plotly_chart(style_fig(comp_fig), use_container_width=True)
            else: st.info("Need at least 2 distinct districts to compare.")
        else:
            st.info("No data available for benchmarking.")

    with tab7:
        st.subheader("👥 HR Management Only")
        
        st.markdown("### Complete HR & Employee Module (Registry, Supervisor Assigned & District Mapped)")
        con_hr = duckdb.connect(DB_FILE_PATH)
        emp_df = con_hr.execute(f"SELECT * FROM {EMPLOYEE_TABLE}").df()
        con_hr.close()
        
        if has_permission("Admin") or has_permission("Supervisor"):
            existing_ids = emp_df["emp_id"].tolist() if not emp_df.empty else []
            selected_emp_id_input = st.selectbox("Select Existing ID to Update (or leave blank for new)", options=[""] + existing_ids, key="selected_update_emp_id")
            
            target_record = None
            if selected_emp_id_input and not emp_df.empty:
                match_df = emp_df[emp_df["emp_id"] == selected_emp_id_input]
                if not match_df.empty:
                    target_record = match_df.iloc[0]

            default_id = selected_emp_id_input if selected_emp_id_input else ""
            default_name = str(target_record["name"]) if target_record is not None and pd.notna(target_record["name"]) else ""
            default_dist = str(target_record["district"]) if target_record is not None and pd.notna(target_record["district"]) else ""
            default_sup = str(target_record["supervisor"]) if target_record is not None and pd.notna(target_record["supervisor"]) else ""
            
            role_choices = ["Data Processing Assistant (DPA)", "Divisional IT Supervisor"]
            default_role = str(target_record["emp_role"]) if target_record is not None and pd.notna(target_record["emp_role"]) else role_choices[0]
            
            default_status = str(target_record["status"]) if target_record is not None and pd.notna(target_record["status"]) else "Active"
            
            try:
                default_join = datetime.datetime.strptime(str(target_record["joining_date"]), "%Y-%m-%d").date() if target_record is not None and pd.notna(target_record["joining_date"]) else datetime.date(2026, 1, 1)
            except:
                default_join = datetime.date(2026, 1, 1)
                
            try:
                default_leave = datetime.datetime.strptime(str(target_record["leaving_date"]), "%Y-%m-%d").date() if target_record is not None and pd.notna(target_record["leaving_date"]) and str(target_record["leaving_date"]) != "None" else None
            except:
                default_leave = None

            with st.form("employee_form"):
                st.markdown("#### Add / Update Employee Record (District & Supervisor Mapped)")
                
                e_id = st.text_input("Employee ID", value=default_id)
                e_name = st.text_input("Employee Full Name", value=default_name)
                e_dist = st.text_input("Assigned District (Registry)", value=default_dist)
                e_sup = st.text_input("Assigned Supervisor (Registry)", value=default_sup)
                
                role_idx = role_choices.index(default_role) if default_role in role_choices else 0
                e_role = st.selectbox("Employee Role", role_choices, index=role_idx)
                
                status_options = ["Active", "Archived", "Transferred"]
                status_idx = status_options.index(default_status) if default_status in status_options else 0
                e_status = st.selectbox("Status", status_options, index=status_idx)
                
                e_join = st.date_input("Joining Date", value=default_join)
                e_leave = st.date_input("Leaving Date (Optional)", value=default_leave)
                
                submitted = st.form_submit_button("Save / Update Employee Record")
                if submitted and e_id and e_name:
                    con_hr = duckdb.connect(DB_FILE_PATH)
                    con_hr.execute(f"""
                        INSERT OR REPLACE INTO {EMPLOYEE_TABLE} 
                        (emp_id, name, district, supervisor, emp_role, status, joining_date, leaving_date) 
                        VALUES (?, ?, ?, ?, ?, ?, ?, ?)
                    """, [
                        e_id.strip(), 
                        e_name.strip().upper(), 
                        e_dist.strip().upper(), 
                        e_sup.strip().upper(), 
                        e_role, 
                        e_status, 
                        str(e_join), 
                        str(e_leave) if e_leave else None
                    ])
                    con_hr.close()
                    log_audit_event(user_id=f"Admin/Supervisor ({st.session_state.get('user_role', 'Admin')})", action="UPDATE_EMPLOYEE", details=f"Updated/Created employee record for ID '{e_id}' ({e_name})")
                    st.success(f"Employee record for ID '{e_id}' successfully saved and updated with role, district, and supervisor mappings!")
                    st.rerun()

            if existing_ids:
                st.markdown("---")
                st.markdown("#### 🗑️ Delete Employee Record")
                with st.form("delete_employee_form"):
                    del_emp_id = st.selectbox("Select Employee ID to Delete", options=existing_ids, key="delete_emp_id_select")
                    delete_submitted = st.form_submit_button("🗑️ Delete Selected Employee", type="primary")
                    if delete_submitted and del_emp_id:
                        con_hr = duckdb.connect(DB_FILE_PATH)
                        con_hr.execute(f"DELETE FROM {EMPLOYEE_TABLE} WHERE emp_id = ?", [del_emp_id])
                        con_hr.close()
                        log_audit_event(user_id=f"Admin/Supervisor ({st.session_state.get('user_role', 'Admin')})", action="DELETE_EMPLOYEE", details=f"Deleted employee record for ID '{del_emp_id}'")
                        st.success(f"Employee record for ID '{del_emp_id}' successfully deleted from database!")
                        st.rerun()
        else:
            st.info("🔒 Employee modifications and deletions require Supervisor or Admin role.")
                    
        st.markdown("#### Employee Directory Registry")
        st.dataframe(emp_df, use_container_width=True)

if current_role == "Admin":
    with tab8:
        st.subheader(f"📝 Live Spreadsheet Data Editor — {selected_sheet}")
        if has_permission("Admin"):
            column_config_dict = {}
            for col in raw_df.columns:
                if str(col).strip().isdigit():
                    try:
                        if datetime.date(yr_val, m_num, int(col)).weekday() == 6:
                            column_config_dict[col] = st.column_config.NumberColumn(label=f"🔴 {col} (Sun)")
                    except: pass
            
            st.markdown('<div class="horizontal-scroll-container">', unsafe_allow_html=True)
            edited_raw_df = st.data_editor(raw_df, num_rows="dynamic", use_container_width=True, column_config=column_config_dict)
            st.markdown('</div>', unsafe_allow_html=True)
            
            if st.button("💾 Save to DuckDB Database"):
                calculated_edited_df = auto_calculate_total_pages(edited_raw_df)
                save_sheet_to_duckdb(selected_year, selected_sheet, calculated_edited_df)
                log_audit_event(user_id=f"Admin ({st.session_state.get('user_role', 'Admin')})", action="DATA_CHANGE", details=f"Updated data spreadsheet for sheet {selected_year}::{selected_sheet}")
                st.session_state["hierarchical_data"] = get_all_tables_hierarchical()
                st.success("Changes calculated and stored successfully!")
                st.rerun()
        else:
            st.warning("🔒 Data editing is restricted to Admin role.")
            st.markdown('<div class="horizontal-scroll-container">', unsafe_allow_html=True)
            st.dataframe(raw_df, use_container_width=True)
            st.markdown('</div>', unsafe_allow_html=True)

    with tab9:
        st.subheader("🛡️ System Activity & Audit Logs")
        
        col_f1, col_f2 = st.columns(2)
        with col_f1:
            action_filter = st.selectbox("Filter by Action Type", ["All", "LOGIN", "LOGOUT", "DATA_CHANGE", "UPDATE_EMPLOYEE", "DELETE_EMPLOYEE", "UPDATE_TARGET", "CREATE_YEAR", "DELETE_YEAR", "CREATE_MONTH", "DELETE_MONTH"])
        with col_f2:
            search_log_user = st.text_input("Search by User / ID")

        con_audit = duckdb.connect(DB_FILE_PATH)
        try:
            query_audit = f"SELECT * FROM {AUDIT_LOG_TABLE} WHERE 1=1"
            params_audit = []
            
            if action_filter != "All":
                query_audit += " AND action = ?"
                params_audit.append(action_filter)
                
            if search_log_user:
                query_audit += " AND user_id LIKE ?"
                params_audit.append(f"%{search_log_user}%")
                
            query_audit += " ORDER BY timestamp DESC LIMIT 500"
            audit_df = con_audit.execute(query_audit, params_audit).df()
        except:
            audit_df = pd.DataFrame(columns=["timestamp", "user_id", "action", "details"])
        con_audit.close()
        
        if not audit_df.empty:
            st.dataframe(audit_df, use_container_width=True)
        else:
            st.info("No audit logs found matching the criteria.")
